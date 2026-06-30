import json
from fastapi.testclient import TestClient
from app import app

client = TestClient(app)


def _files(sample_template_dir):
    return {"file": ("base.pptx", (sample_template_dir / "base.pptx").read_bytes())}


def test_health():
    assert client.get("/health").json() == {"ok": True}


def test_extract_shapes(sample_template_dir):
    r = client.post("/extract-shapes", files=_files(sample_template_dir))
    assert r.status_code == 200
    assert len(r.json()["slides"]) == 4


def test_render_deck_ok(sample_template_dir, sample_manifest):
    deck = {"slides": [{"slide_type": "title", "slots": {"title": "Hi", "subtitle": "Yo"}}]}
    r = client.post("/render-deck", files=_files(sample_template_dir),
                    data={"manifest": json.dumps(sample_manifest), "deck_spec": json.dumps(deck)})
    assert r.status_code == 200
    assert r.content[:2] == b"PK"  # zip/pptx magic


def test_render_deck_rejects(sample_template_dir, sample_manifest):
    # A wrong-typed value is a fatal validation error -> 422. (Text overflow is
    # no longer fatal: it shrinks then sentence-truncates with a non-fatal
    # X-Overflow-Warnings header, so it is not a rejection case.)
    deck = {"slides": [{"slide_type": "title", "slots": {"title": 123}}]}
    r = client.post("/render-deck", files=_files(sample_template_dir),
                    data={"manifest": json.dumps(sample_manifest), "deck_spec": json.dumps(deck)})
    assert r.status_code == 422
    assert r.json()["validation"][0]["code"] == "wrong_type"


def test_move_shape(sample_template_dir):
    from pptx import Presentation
    sid = Presentation(str(sample_template_dir / "base.pptx")).slides[0].shapes[0].shape_id
    r = client.post("/move-shape", files=_files(sample_template_dir),
                    data={"shape_id": str(sid),
                          "bbox_pct": json.dumps({"x": 10, "y": 10, "w": 40, "h": 20})})
    assert r.status_code == 200
    assert r.content[:2] == b"PK"


def test_move_shapes_endpoint():
    import io
    from pptx import Presentation

    def _two_slide_deck() -> bytes:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for _ in range(2):
            s = prs.slides.add_slide(blank)
            s.shapes.add_textbox(0, 0, 914400, 914400)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    data = _two_slide_deck()
    sid = Presentation(io.BytesIO(data)).slides[1].shapes[0].shape_id
    moves = [{"slide_index": 1, "shape_id": sid,
              "bbox_pct": {"x": 50, "y": 50, "w": 20, "h": 10}}]
    r = client.post("/move-shapes",
                    files={"file": ("d.pptx", data)},
                    data={"moves": json.dumps(moves)})
    assert r.status_code == 200
    assert r.content[:2] == b"PK"


def test_validate_deck_ok(sample_template_dir, sample_manifest):
    deck_spec = {"slides": [{"slide_type": "title", "slots": {"title": "Hi", "subtitle": "Yo"}}]}
    r = client.post("/validate-deck", files=_files(sample_template_dir),
                    data={"manifest": json.dumps(sample_manifest),
                          "deck_spec": json.dumps(deck_spec)})
    assert r.status_code == 200
    body = r.json()
    assert "errors" in body and "warnings" in body
    assert body["errors"] == []


def test_validate_deck_reports_errors(sample_template_dir, sample_manifest):
    # An invalid deck (unknown slide_type) -> 200 with errors in the body.
    r = client.post("/validate-deck", files=_files(sample_template_dir),
                    data={"manifest": json.dumps(sample_manifest),
                          "deck_spec": json.dumps({"slides": [{"slide_type": "nope", "slots": {}}]})})
    assert r.status_code == 200
    assert len(r.json()["errors"]) >= 1


def test_render_base_previews_timeout_returns_note(sample_template_dir, monkeypatch):
    import app as app_mod
    from pptx_mcp.preview import PreviewTimeout
    monkeypatch.setattr(app_mod, "libreoffice_available", lambda: True)
    def boom(_data):
        raise PreviewTimeout("soffice timed out")
    monkeypatch.setattr(app_mod, "preview", boom)
    r = client.post("/render-base-previews", files=_files(sample_template_dir))
    assert r.status_code == 200
    assert r.json() == {"previews": [], "note": "preview timed out"}


def test_render_preview_timeout_returns_note(sample_template_dir, sample_manifest, monkeypatch):
    import app as app_mod
    from pptx_mcp.preview import PreviewTimeout
    monkeypatch.setattr(app_mod, "libreoffice_available", lambda: True)
    def boom(_data):
        raise PreviewTimeout("soffice timed out")
    monkeypatch.setattr(app_mod, "preview", boom)
    deck = {"slides": [{"slide_type": "title", "slots": {"title": "Hi", "subtitle": "Yo"}}]}
    r = client.post("/render-preview", files=_files(sample_template_dir),
                    data={"manifest": json.dumps(sample_manifest), "deck_spec": json.dumps(deck)})
    assert r.status_code == 200
    assert r.json() == {"validation": [], "previews": [], "note": "preview timed out"}
