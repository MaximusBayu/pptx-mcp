import io

from pptx.util import Pt

from .assembler import find_shape
from .fit import assess_text
from .models import Slot

_BASE_PT = 24.0
_SHRINK_STEP = 4.0


def fill_slot(slide, slot: Slot, value) -> None:
    shape = find_shape(slide, slot.shape_id)
    if slot.type == "text":
        _fill_text(shape, slot, value)
    elif slot.type == "table":
        _fill_table(shape, value)
    elif slot.type == "image":
        _fill_image(slide, shape, value)


def _fill_text(shape, slot: Slot, value: str) -> None:
    tf = shape.text_frame
    tf.text = value
    decision, _ = assess_text(value, slot.constraints)
    if decision == "shrink":
        floor = slot.constraints.shrink_floor_pt or 12.0
        new_pt = max(floor, _BASE_PT - _SHRINK_STEP)
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(new_pt)


def _fill_table(shape, rows: list[list]) -> None:
    table = shape.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if r < len(table.rows) and c < len(table.columns):
                table.cell(r, c).text = str(val)


def _fill_image(slide, shape, value) -> None:
    data = value if isinstance(value, bytes) else open(value, "rb").read()
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(io.BytesIO(data), left, top, width, height)
