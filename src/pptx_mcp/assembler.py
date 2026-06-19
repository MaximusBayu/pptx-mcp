import copy

from pptx import Presentation

from .models import Template

# Namespace URI for Office Open XML relationships
_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def find_shape(slide, shape_id: int):
    """Return the shape on *slide* whose shape_id matches, or raise KeyError."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise KeyError(f"shape_id {shape_id} not found on slide")


def _duplicate_slide(prs: Presentation, src_index: int):
    """Append a copy of slide *src_index* (within the same package) and return it.

    Deep-copies every shape element from the source slide's spTree into the new
    slide.  For picture shapes the ``<a:blip r:embed="rIdN">`` attribute references
    a relationship id in the *source* slide's part; after copying, those ids may
    or may not match ids in the *destination* part.  This function:

    1. Copies each image relationship from the source part into the dest part via
       ``Part.relate_to``, which returns the (possibly different) rId assigned in
       the dest part.
    2. Rewrites every ``r:embed`` and ``r:link`` attribute in the copied XML tree
       from the old source rId to the new dest rId.
    """
    source = prs.slides[src_index]
    layout = source.slide_layout
    dest = prs.slides.add_slide(layout)

    # Remove any placeholder shapes that add_slide inserted from the layout.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # Deep-copy each source shape element into the dest spTree.
    for shp in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    # Copy image relationships and record old-rId -> new-rId mapping.
    old_to_new_rid: dict[str, str] = {}
    for rid, rel in source.part.rels.items():
        if "image" in rel.reltype:
            new_rid = dest.part.relate_to(rel._target, rel.reltype)
            old_to_new_rid[rid] = new_rid

    # Rewrite r:embed / r:link attributes in the copied XML to the dest rIds.
    if old_to_new_rid:
        embed_attr = "{%s}embed" % _REL_NS
        link_attr = "{%s}link" % _REL_NS
        for elem in dest.shapes._spTree.iter():
            old_embed = elem.get(embed_attr)
            if old_embed and old_embed in old_to_new_rid:
                elem.set(embed_attr, old_to_new_rid[old_embed])
            old_link = elem.get(link_attr)
            if old_link and old_link in old_to_new_rid:
                elem.set(link_attr, old_to_new_rid[old_link])

    return dest


def assemble(order: list[int], template: Template) -> Presentation:
    """Return a Presentation whose slides are copies of the template's base slides
    at the given *source_slide_index* values (in *order*), with the original base
    slides removed.

    All shape types (text, tables, pictures) survive a save -> reopen round-trip.
    """
    prs = Presentation(template.pptx_path)
    original_count = len(prs.slides)

    for src_index in order:
        _duplicate_slide(prs, src_index)

    # Remove the original base slides (indices 0 .. original_count-1).
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for sid in slide_ids[:original_count]:
        xml_slides.remove(sid)

    return prs
