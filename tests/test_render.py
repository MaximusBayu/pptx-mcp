import io
import pytest
from pptx import Presentation
from pptx_mcp.template import load_template
from pptx_mcp.render import render, RenderRejected


def _deck():
    return {"slides": [
        {"slide_type": "title", "slots": {"title": "Acme", "subtitle": "Q3"}},
        {"slide_type": "bullet", "slots": {"heading": "Plan", "body": "a\nb"}},
    ]}


def test_render_produces_valid_pptx(sample_template_dir):
    tpl = load_template(sample_template_dir)
    data = render(_deck(), tpl)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 2
    texts = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
    assert "Acme" in texts


def test_render_rejects_invalid(sample_template_dir):
    tpl = load_template(sample_template_dir)
    bad = {"slides": [{"slide_type": "title", "slots": {"title": "x" * 200}}]}
    with pytest.raises(RenderRejected) as ei:
        render(bad, tpl)
    assert ei.value.errors[0].code == "text_overflow"
