import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _guess_type(shape) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    return "text"


def _pct(value, total) -> float:
    return round(min(100.0, max(0.0, 100.0 * value / total)), 3) if total else 0.0


def extract_shapes(pptx_bytes: bytes) -> dict:
    prs = Presentation(io.BytesIO(pptx_bytes))
    sw, sh = prs.slide_width, prs.slide_height
    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = []
        for shp in slide.shapes:
            x, y, w, h = (shp.left if shp.left is not None else 0), (shp.top if shp.top is not None else 0), (shp.width if shp.width is not None else 0), (shp.height if shp.height is not None else 0)
            shapes.append({
                "shape_id": shp.shape_id, "name": shp.name or "",
                "type": _guess_type(shp),
                "x": x, "y": y, "w": w, "h": h,
                "bbox_pct": {"x": _pct(x, sw), "y": _pct(y, sh),
                             "w": _pct(w, sw), "h": _pct(h, sh)},
            })
        slides.append({"index": i, "width_emu": sw, "height_emu": sh, "shapes": shapes})
    return {"slides": slides}
