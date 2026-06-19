import argparse
import json
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _slug(name: str | None, fallback: str) -> str:
    s = re.sub(r"[^a-z0-9]+", "_", (name or "").lower()).strip("_")
    return s or fallback


def _guess_type(shape) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    return "text"


def scaffold_manifest(pptx_path: str) -> dict:
    prs = Presentation(pptx_path)
    slide_types = []
    for i, slide in enumerate(prs.slides):
        slots = []
        for j, shape in enumerate(slide.shapes):
            stype = _guess_type(shape)
            if stype == "text" and not shape.has_text_frame:
                continue
            slots.append({
                "id": _slug(shape.name, f"slot_{j}"),
                "name": shape.name or f"Slot {j}",
                "type": stype,
                "target": {"shape_id": shape.shape_id},
                "required": stype != "image",
                "default": None,
                "constraints": {},
            })
        slide_types.append({
            "id": f"slide_{i}", "name": f"Slide {i}", "description": "",
            "source_slide_index": i, "slots": slots,
        })
    return {"template": {"id": "TODO", "name": "TODO", "description": ""},
            "slide_types": slide_types}


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(prog="pptx-mcp")
    sub = parser.add_subparsers(dest="cmd", required=True)
    init = sub.add_parser("init-template")
    init.add_argument("pptx")
    init.add_argument("-o", "--out", required=True)
    args = parser.parse_args(argv)
    if args.cmd == "init-template":
        manifest = scaffold_manifest(args.pptx)
        with open(args.out, "w", encoding="utf-8") as f:
            json.dump(manifest, f, indent=2)
        return 0
    return 1
