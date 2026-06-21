import io

from pptx import Presentation


def move_shape(pptx_bytes: bytes, shape_id: int, bbox_pct: dict,
               slide_index: int | None = None) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))
    sw, sh = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides):
        if slide_index is not None and i != slide_index:
            continue
        for shp in slide.shapes:
            if shp.shape_id == shape_id:
                shp.left = int(sw * bbox_pct["x"] / 100.0)
                shp.top = int(sh * bbox_pct["y"] / 100.0)
                shp.width = int(sw * bbox_pct["w"] / 100.0)
                shp.height = int(sh * bbox_pct["h"] / 100.0)
                buf = io.BytesIO()
                prs.save(buf)
                return buf.getvalue()
    raise KeyError(f"shape_id {shape_id} not found")


def move_shapes(pptx_bytes: bytes, moves: list[dict]) -> bytes:
    """Apply many moves in one pass. Each move is
    {slide_index, shape_id, bbox_pct:{x,y,w,h}} with bbox in slide-percent."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    sw, sh = prs.slide_width, prs.slide_height
    slides = list(prs.slides)
    for m in moves:
        si = m["slide_index"]
        if not (0 <= si < len(slides)):
            raise KeyError(f"slide_index {si} out of range")
        shp = next((s for s in slides[si].shapes if s.shape_id == m["shape_id"]), None)
        if shp is None:
            raise KeyError(f"shape_id {m['shape_id']} not found on slide {si}")
        b = m["bbox_pct"]
        shp.left = int(sw * b["x"] / 100.0)
        shp.top = int(sh * b["y"] / 100.0)
        shp.width = int(sw * b["w"] / 100.0)
        shp.height = int(sh * b["h"] / 100.0)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
