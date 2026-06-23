import base64
import io
import logging
import math
import urllib.request

from PIL import Image
from pptx.util import Emu, Length, Pt

from .assembler import find_shape
from .autodetect import LINE_H, estimate_max_chars
from .models import Slot, SlotError
from .tablefit import MIN_COL_FRAC, MIN_ROW_FRAC, redistribute
from .textfit import fit_text, truncate_to_sentence

_BASE_PT = 24.0
_MIN_PT = 12.0
_CELL_SHRINK_PT = 4.0
_IMG_MAX_BYTES = 20 * 1024 * 1024


def load_image_bytes(value) -> bytes:
    """Resolve an image slot value to raw bytes.

    Accepts raw bytes, a data: URL (base64), an http(s) URL, or a local file
    path (the last is for engine-local/test use). Remote fetches are capped in
    size and restricted to http(s); the engine runs on input proxied by the
    web API, but be aware a URL value triggers a server-side fetch.
    """
    if isinstance(value, bytes):
        return value
    if not isinstance(value, str) or not value:
        raise ValueError("image value must be a non-empty str or bytes")
    if value.startswith("data:"):
        _, _, b64 = value.partition(",")
        return base64.b64decode(b64)
    if value.startswith(("http://", "https://")):
        req = urllib.request.Request(value, headers={"User-Agent": "pptx-mcp"})
        with urllib.request.urlopen(req, timeout=15) as resp:  # noqa: S310 (scheme checked)
            data = resp.read(_IMG_MAX_BYTES + 1)
        if len(data) > _IMG_MAX_BYTES:
            raise ValueError("image exceeds size limit")
        return data
    with open(value, "rb") as fh:
        return fh.read()


def fill_slot(slide, slot: Slot, value) -> list[SlotError]:
    shape = find_shape(slide, slot.shape_id)
    if slot.type == "text":
        return _fill_text(shape, slot, value)
    if slot.type == "table":
        return _fill_table(shape, value)
    if slot.type == "image":
        _fill_image(slide, shape, value, slot.constraints.fit)
    return []


def clear_slot(slide, slot: Slot) -> None:
    """Remove the template's sample content for an unfilled slot.

    text  -> empty the text frame
    image -> remove the placeholder/picture shape from the slide
    table -> blank every cell
    A slot whose shape is not on the slide is a silent no-op.
    """
    try:
        shape = find_shape(slide, slot.shape_id)
    except KeyError:
        return
    if slot.type == "text":
        _clear_text(shape)
    elif slot.type == "image":
        _clear_image(slide, shape)
    elif slot.type == "table":
        _clear_table(shape)


def _resolve_spacing(p0, orig_pt) -> float:
    # python-pptx line_spacing is None, a float multiple, or a Length (fixed
    # distance). Resolve to a multiple. Length subclasses int, so check it first.
    ls = p0.line_spacing if p0 is not None else None
    if isinstance(ls, Length):
        return max(0.5, min(3.0, ls.pt / orig_pt))
    if isinstance(ls, (int, float)) and ls > 0:
        return float(ls)
    return LINE_H


def _fill_text(shape, slot: Slot, value: str) -> list[SlotError]:
    warnings: list[SlotError] = []
    tf = shape.text_frame
    tf.word_wrap = True

    # Preserve the template's styling: keep the first paragraph (its alignment)
    # and write into its first run (its font family, size, bold/italic, color).
    p0 = tf.paragraphs[0] if tf.paragraphs else None
    r0 = p0.runs[0] if (p0 is not None and p0.runs) else None
    orig_pt = r0.font.size.pt if (r0 is not None and r0.font.size is not None) else _BASE_PT

    base_spacing = _resolve_spacing(p0, orig_pt)
    floor_pt = slot.constraints.shrink_floor_pt or _MIN_PT
    res = fit_text(value, shape.width or 0, shape.height or 0, orig_pt, floor_pt, base_spacing)
    value = res.value
    dropped = res.dropped

    # Preserve the existing hard max_chars cap on top of the geometric fit.
    max_chars = slot.constraints.max_chars
    if max_chars is not None and len(value) > max_chars:
        value, extra = truncate_to_sentence(value, max_chars)
        dropped = (dropped + extra) if dropped else extra

    if dropped:
        warnings.append(SlotError(0, slot.id, "text_truncated",
                                  f"dropped {len(dropped)} chars to fit"))

    if r0 is not None:
        # Write into the existing run; drop extra runs and paragraphs so the
        # template's formatting on r0/p0 is what remains.
        r0.text = value
        for extra_run in p0.runs[1:]:
            extra_run._r.getparent().remove(extra_run._r)
        for extra_para in tf.paragraphs[1:]:
            extra_para._p.getparent().remove(extra_para._p)
        r0.font.size = Pt(res.font_pt)
        p0.line_spacing = res.line_spacing
    else:
        # No run to inherit from (empty box) — fall back to plain text.
        tf.text = value
        for para in tf.paragraphs:
            para.line_spacing = res.line_spacing
            for run in para.runs:
                run.font.size = Pt(res.font_pt)
    return warnings


def _cell_font_pt(cell) -> float:
    tf = cell.text_frame
    p0 = tf.paragraphs[0] if tf.paragraphs else None
    r0 = p0.runs[0] if (p0 is not None and p0.runs) else None
    if r0 is not None and r0.font.size is not None:
        return r0.font.size.pt
    return _BASE_PT


def _any_cell_overflows(table, rows, col_w, row_h) -> bool:
    for r, row in enumerate(rows):
        if r >= len(table.rows):
            continue
        for c, val in enumerate(row):
            if c >= len(table.columns):
                continue
            cap, _ = estimate_max_chars(col_w[c], row_h[r], _cell_font_pt(table.cell(r, c)))
            if len(str(val)) > cap:
                return True
    return False


def _resize_columns(table, rows, col_w) -> list[int]:
    ncols = len(col_w)
    demands = [0.0] * ncols
    for row in rows:
        for c, val in enumerate(row):
            if c < ncols:
                demands[c] = max(demands[c], float(len(str(val))))
    total = sum(col_w)
    return redistribute(demands, total, int(MIN_COL_FRAC * total))


def _resize_rows(table, rows, col_w, row_h) -> list[int]:
    nrows = len(row_h)
    demands = [0.0] * nrows
    for r, row in enumerate(rows):
        if r >= nrows:
            continue
        for c, val in enumerate(row):
            if c >= len(col_w):
                continue
            cap, lines = estimate_max_chars(col_w[c], row_h[r], _cell_font_pt(table.cell(r, c)))
            cpl = max(1, cap // lines)  # chars per line at the cell's font
            lines_needed = math.ceil(len(str(val)) / cpl)
            demands[r] = max(demands[r], float(lines_needed))
    total = sum(row_h)
    return redistribute(demands, total, int(MIN_ROW_FRAC * total))


def _fit_cell(cell, value, width_emu, height_emu, slot_id) -> list[SlotError]:
    warnings: list[SlotError] = []
    tf = cell.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0] if tf.paragraphs else None
    r0 = p0.runs[0] if (p0 is not None and p0.runs) else None
    orig_pt = r0.font.size.pt if (r0 is not None and r0.font.size is not None) else _BASE_PT

    capacity, _ = estimate_max_chars(width_emu, height_emu, orig_pt)
    new_pt = orig_pt
    if len(value) > capacity:
        new_pt = max(_MIN_PT, orig_pt - _CELL_SHRINK_PT)
        capacity2 = int(capacity * (orig_pt / new_pt))
        if len(value) > capacity2:
            value, dropped = truncate_to_sentence(value, capacity2)
            if dropped:
                warnings.append(SlotError(0, slot_id, "text_truncated",
                                          f"dropped {len(dropped)} chars to fit cell"))

    shrunk = new_pt < orig_pt

    if r0 is not None:
        r0.text = value
        for extra_run in p0.runs[1:]:
            extra_run._r.getparent().remove(extra_run._r)
        for extra_para in tf.paragraphs[1:]:
            extra_para._p.getparent().remove(extra_para._p)
        if shrunk:
            r0.font.size = Pt(new_pt)
    else:
        cell.text = value
        if shrunk:
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(new_pt)
    return warnings


def _fill_table(shape, rows: list[list]) -> list[SlotError]:
    table = shape.table
    _blank_all_cells(table)          # drop template sample rows before filling
    warnings: list[SlotError] = []
    col_w = [table.columns[c].width for c in range(len(table.columns))]
    row_h = [table.rows[r].height for r in range(len(table.rows))]

    # Auto-size only when some provided cell overflows at its base font, so a
    # table that already fits keeps its original column/row sizes.
    if _any_cell_overflows(table, rows, col_w, row_h):
        col_w = _resize_columns(table, rows, col_w)
        for c, w in enumerate(col_w):
            table.columns[c].width = Emu(w)
        row_h = _resize_rows(table, rows, col_w, row_h)  # demand uses new widths
        for r, h in enumerate(row_h):
            table.rows[r].height = Emu(h)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if r < len(table.rows) and c < len(table.columns):
                warnings.extend(_fit_cell(table.cell(r, c), str(val),
                                          col_w[c], row_h[r], f"cell[{r},{c}]"))
    return warnings


def _fill_image(slide, shape, value, fit: str | None = None) -> None:
    data = load_image_bytes(value)
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)

    new_left, new_top, new_w, new_h = left, top, width, height
    # "contain" (default): scale to fit inside the box, preserve aspect, center.
    # "cover" is deferred -> treated as contain for now.
    if width and height:
        try:
            iw, ih = Image.open(io.BytesIO(data)).size
        except Exception:
            iw = ih = 0
        if iw > 0 and ih > 0:
            box_ar = width / height
            img_ar = iw / ih
            if img_ar > box_ar:
                new_w = width
                new_h = round(width / img_ar)
            else:
                new_h = height
                new_w = round(height * img_ar)
            new_left = left + (width - new_w) // 2
            new_top = top + (height - new_h) // 2

    try:
        slide.shapes.add_picture(io.BytesIO(data), new_left, new_top, new_w, new_h)
    except Exception as exc:
        # Fallback: create a placeholder image if the input is invalid
        logging.getLogger(__name__).warning("image slot fill failed; inserting placeholder at box rect: %s", exc)
        buf = io.BytesIO()
        Image.new("RGB", (1, 1), (200, 200, 200)).save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, new_left, new_top, new_w, new_h)


def _clear_text(shape) -> None:
    shape.text_frame.clear()


def _clear_image(slide, shape) -> None:
    shape._element.getparent().remove(shape._element)


def _blank_all_cells(table) -> None:
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            table.cell(r, c).text = ""


def _clear_table(shape) -> None:
    _blank_all_cells(shape.table)
