from pptx import Presentation

from .models import Constraints, Slot, SlideType, Template

_VALID_TYPES = {"text", "table", "image"}


class ManifestError(Exception):
    pass


def _constraints(d: dict | None) -> Constraints:
    d = d or {}
    return Constraints(
        max_chars=d.get("max_chars"), max_lines=d.get("max_lines"),
        shrink_floor_pt=d.get("shrink_floor_pt"),
        max_rows=d.get("max_rows"), max_cols=d.get("max_cols"),
        fit=d.get("fit"),
    )


def parse_manifest(data: dict, pptx_path: str) -> Template:
    try:
        meta = data["template"]
        slide_types = []
        for st in data["slide_types"]:
            slots = []
            for s in st["slots"]:
                if s["type"] not in _VALID_TYPES:
                    raise ManifestError(f"invalid slot type: {s['type']}")
                slots.append(Slot(
                    id=s["id"], name=s["name"], type=s["type"],
                    shape_id=s["target"]["shape_id"],
                    required=s.get("required", False),
                    default=s.get("default"),
                    constraints=_constraints(s.get("constraints")),
                ))
            slide_types.append(SlideType(
                id=st["id"], name=st["name"], description=st.get("description", ""),
                source_slide_index=st["source_slide_index"], slots=slots,
            ))
        return Template(id=meta["id"], name=meta["name"],
                        description=meta.get("description", ""),
                        slide_types=slide_types, pptx_path=pptx_path)
    except KeyError as e:
        raise ManifestError(f"missing manifest key: {e}") from e


def validate_against_pptx(template: Template) -> None:
    prs = Presentation(template.pptx_path)
    n = len(prs.slides)
    for st in template.slide_types:
        if not (0 <= st.source_slide_index < n):
            raise ManifestError(
                f"slide_type {st.id}: source_slide_index {st.source_slide_index} out of range (0..{n-1})")
        slide = prs.slides[st.source_slide_index]
        present = {sh.shape_id for sh in slide.shapes}
        for slot in st.slots:
            if slot.shape_id not in present:
                raise ManifestError(
                    f"slide_type {st.id} slot {slot.id}: shape_id {slot.shape_id} "
                    f"not found on slide {st.source_slide_index}")
