from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

_TEXT_SAMPLE_MAX = 200


def component_type(shp) -> str:
    if getattr(shp, "has_table", False):
        return "table"
    if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if getattr(shp, "has_text_frame", False):
        return "text"
    return "other"


def _is_multiline(shp) -> bool:
    if not getattr(shp, "has_text_frame", False):
        return False
    paras = shp.text_frame.paragraphs
    if sum(1 for p in paras if (p.text or "").strip()) > 1:
        return True
    p0 = paras[0] if paras else None
    if p0 is not None:
        pPr = p0._p.find(qn("a:pPr"))
        if pPr is not None and (pPr.find(qn("a:buChar")) is not None
                                or pPr.find(qn("a:buAutoNum")) is not None):
            return True
    return False


def _hint(ctype: str, multiline: bool) -> str:
    if ctype == "text":
        return ("bullet list — pass content as an array of strings, one per bullet"
                if multiline else "single text — pass a string")
    if ctype == "table":
        return "pass rows as list[list]"
    if ctype == "image":
        return "pass a URL or base64 string"
    return "decorative — placed verbatim, no content"


def _pct(value, total) -> float:
    return round(min(100.0, max(0.0, 100.0 * value / total)), 3) if total else 0.0


def _sample_text(shp) -> str:
    if not getattr(shp, "has_text_frame", False):
        return ""
    t = (shp.text_frame.text or "").strip()
    return t[: _TEXT_SAMPLE_MAX - 1] + "…" if len(t) > _TEXT_SAMPLE_MAX else t


def _hex_or_none(color):
    # color.rgb raises (TypeError/AttributeError) for theme/inherited colors.
    try:
        if color is not None and color.type is not None:
            return str(color.rgb)
    except (TypeError, AttributeError):
        pass
    return None


def _shape_style(shp) -> dict:
    style = {"font_name": None, "font_pt": None, "font_color": None, "fill_color": None}
    if getattr(shp, "has_text_frame", False):
        paras = shp.text_frame.paragraphs
        runs = paras[0].runs if paras else []
        if runs:
            f = runs[0].font
            style["font_name"] = f.name
            style["font_pt"] = f.size.pt if f.size is not None else None
            style["font_color"] = _hex_or_none(f.color)
    try:
        fill = shp.fill
        if fill.type is not None:
            style["fill_color"] = _hex_or_none(fill.fore_color)
    except (TypeError, AttributeError, ValueError):
        pass
    return style


def _component_dict(shp, slide_index, sw, sh, slot_id) -> dict:
    x = shp.left or 0
    y = shp.top or 0
    w = shp.width or 0
    h = shp.height or 0
    ctype = component_type(shp)
    multiline = _is_multiline(shp)
    return {
        "component_id": f"{slide_index}:{shp.shape_id}",
        "source_slide": slide_index,
        "type": ctype,
        "multiline": multiline,
        "hint": _hint(ctype, multiline),
        "fillable": slot_id is not None,
        "slot_id": slot_id,
        "name": shp.name or "",
        "geometry": {
            "bbox_pct": {"x": _pct(x, sw), "y": _pct(y, sh),
                         "w": _pct(w, sw), "h": _pct(h, sh)},
            "width_emu": int(w), "height_emu": int(h),
        },
        "style": _shape_style(shp),
        "text": _sample_text(shp),
    }


def get_catalog(template) -> dict:
    prs = Presentation(template.pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    fillable = {(st.source_slide_index, s.shape_id): s.id
                for st in template.slide_types for s in st.slots}
    components = []
    for i, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            slot_id = fillable.get((i, shp.shape_id))
            components.append(_component_dict(shp, i, sw, sh, slot_id))
    return {"id": template.id, "name": template.name,
            "description": template.description, "components": components}
