import io
import re
from collections import Counter
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .shapes import _guess_type, _pct

TAU = 0.5
GLYPH_W = 0.5
LINE_H = 1.2
DEFAULT_FONT_PT = 18.0
EMU_PER_PT = 12700

_DECO_TYPES = {MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.LINE}
_MIN_AREA_PCT = 1.0
_MIN_DIM_PCT = 0.5
# A tiny box holding only 1-2 characters is almost always a sequence number,
# bullet, or label marker (e.g. the "03"/"04" agenda numbers in real decks),
# not a fillable content slot. Penalize it so it drops below the candidate
# threshold while genuine short fields (a "2026" date, ~4 chars) survive.
_MICRO_AREA_PCT = 30.0
_MICRO_TEXT_LEN = 2


@dataclass
class ShapeAssessment:
    shape_id: int
    name: str
    type: str
    bbox_pct: dict
    confidence: float
    is_candidate: bool
    font_pt: float | None


_EXAMPLE_MAX = 200

_DESC_BY_ID = {
    "title": "Slide title",
    "subtitle": "Subtitle",
    "body": "Body text",
}


def _truncate(text: str, limit: int) -> str:
    text = (text or "").strip()
    if limit <= 0 or len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"


def slot_description(suggested_id: str) -> str:
    """A short human label for a slot, derived from its auto-assigned id."""
    if suggested_id in _DESC_BY_ID:
        return _DESC_BY_ID[suggested_id]
    if suggested_id.startswith("table"):
        return "Table data"
    if suggested_id.startswith("image"):
        return "Image"
    return "Text"


_AGENDA_RE = re.compile(r"agenda|overview|outline|contents|daftar isi", re.I)
_SUMMARY_RE = re.compile(r"summary|ringkasan|executive", re.I)
_FINDING_RE = re.compile(r"finding|temuan|severity|critical|high|medium|low|cwe|cvss", re.I)
_CLOSING_RE = re.compile(r"thank|terima kasih|questions|q&a", re.I)

_KIND_LABEL = {
    "cover": "Cover slide", "agenda": "Agenda slide", "summary": "Summary slide",
    "finding": "Finding slide", "data": "Data slide", "closing": "Closing slide",
    "section": "Section slide", "content": "Content slide",
}


def slide_kind(text_blob: str, has_table: bool, index: int,
               num_text: int, has_subtitle: bool) -> str:
    """Deterministic slide purpose from its dominant text + shape mix."""
    if index == 0 or (has_subtitle and num_text <= 3):
        return "cover"
    if _AGENDA_RE.search(text_blob):
        return "agenda"
    if _SUMMARY_RE.search(text_blob):
        return "summary"
    if _FINDING_RE.search(text_blob):
        return "finding"
    if has_table:
        return "data"
    if _CLOSING_RE.search(text_blob):
        return "closing"
    if num_text <= 1:
        return "section"
    return "content"


def slide_description(kind: str, slot_ids: list[str]) -> str:
    """A templated sentence: what the slide is + which slots to fill."""
    fill = ", ".join(slot_ids) if slot_ids else "no slots"
    repeat = " Repeat per item." if kind == "finding" else ""
    label = _KIND_LABEL.get(kind, "Content slide")
    return f"{label} — fill: {fill}.{repeat}"


def slide_signature(slide: dict) -> tuple:
    """A structural fingerprint of a slide's candidate shapes. Two slides with
    the same signature are the same template pattern repeated (e.g. F1-F4
    findings) and are flagged repeatable."""
    cand = [s for s in slide["shapes"] if s.get("is_candidate")]
    parts = tuple(sorted(
        (s["type"], round((s["bbox_pct"]["w"] * s["bbox_pct"]["h"]) / 10))
        for s in cand
    ))
    sids = tuple(sorted({s["suggested_id"] for s in cand if s["suggested_id"]}))
    return (parts, sids)


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text_frame.text or "").strip()


def _first_font_pt(shape) -> float | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt
    return None


def _has_picture_fill(shape) -> bool:
    # Real templates often place a photo as a picture-filled freeform/autoshape,
    # not a bare PICTURE shape.
    try:
        return shape.fill.type == MSO_FILL.PICTURE
    except Exception:
        return False


def _contains_picture(group) -> bool:
    for sub in getattr(group, "shapes", []):
        try:
            if sub.shape_type == MSO_SHAPE_TYPE.PICTURE or _has_picture_fill(sub):
                return True
            if sub.shape_type == MSO_SHAPE_TYPE.GROUP and _contains_picture(sub):
                return True
        except Exception:
            continue
    return False


def _is_image(shape) -> bool:
    """True for a picture, a picture-filled shape, or a group framing one —
    all are fillable image slots rather than decoration."""
    if _has_picture_fill(shape):
        return True
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return _contains_picture(shape)
    except Exception:
        return False
    return False


def classify_shape(shape, slide_w, slide_h) -> ShapeAssessment:
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    bbox = {"x": _pct(left, slide_w), "y": _pct(top, slide_h),
            "w": _pct(width, slide_w), "h": _pct(height, slide_h)}

    score = 0.5
    text = _shape_text(shape)
    try:
        stype = shape.shape_type
    except Exception:
        stype = None

    # A picture, picture-filled shape, or group framing one is an image slot.
    img = _is_image(shape)
    effective_type = "image" if img else _guess_type(shape)

    # exclude signals
    if stype in _DECO_TYPES and not img:
        score -= 0.5
    if not getattr(shape, "has_text_frame", False) and effective_type == "text":
        score -= 0.3
    if getattr(shape, "has_text_frame", False) and not text:
        score -= 0.25
    area_pct = bbox["w"] * bbox["h"]
    if area_pct < _MIN_AREA_PCT or bbox["w"] < _MIN_DIM_PCT or bbox["h"] < _MIN_DIM_PCT:
        score -= 0.4
    raw_off = (left < 0 or top < 0
               or left + width > slide_w or top + height > slide_h)
    if raw_off:
        score -= 0.4
    # micro label: tiny box with 1-2 chars of text -> sequence number/bullet
    if area_pct < _MICRO_AREA_PCT and 0 < len(text) <= _MICRO_TEXT_LEN:
        score -= 0.5

    # include signals
    if getattr(shape, "is_placeholder", False):
        score += 0.4
    if text and area_pct >= _MIN_AREA_PCT:
        score += 0.3
    if effective_type in ("table", "image") and area_pct >= _MIN_AREA_PCT:
        score += 0.2

    confidence = max(0.0, min(1.0, score))
    return ShapeAssessment(
        shape_id=shape.shape_id, name=shape.name or "",
        type=effective_type, bbox_pct=bbox,
        confidence=round(confidence, 3), is_candidate=confidence >= TAU,
        font_pt=_first_font_pt(shape),
    )


def estimate_max_chars(width_emu, height_emu, font_pt) -> tuple[int, int]:
    pt = font_pt if font_pt and font_pt > 0 else DEFAULT_FONT_PT
    font_emu = pt * EMU_PER_PT
    chars_per_line = max(1, int(width_emu / (font_emu * GLYPH_W)))
    lines = max(1, int(height_emu / (font_emu * LINE_H)))
    return chars_per_line * lines, lines


def autodetect(pptx_bytes: bytes) -> dict:
    prs = Presentation(io.BytesIO(pptx_bytes))
    sw, sh = prs.slide_width, prs.slide_height
    slides = []
    for i, slide in enumerate(prs.slides):
        assessments = [classify_shape(shp, sw, sh) for shp in slide.shapes]
        ids = derive_ids([a for a in assessments if a.is_candidate])
        shape_by_id = {shp.shape_id: shp for shp in slide.shapes}
        shapes = []
        for a in assessments:
            mc = ml = mr = mcols = 0
            if a.is_candidate and a.type == "text":
                shp = shape_by_id[a.shape_id]
                mc, ml = estimate_max_chars(shp.width or 0, shp.height or 0, a.font_pt)
            elif a.is_candidate and a.type == "table":
                shp = shape_by_id[a.shape_id]
                if getattr(shp, "has_table", False):
                    mr = len(shp.table.rows)
                    mcols = len(shp.table.columns)
            shp_obj = shape_by_id[a.shape_id]
            text_val = _truncate(_shape_text(shp_obj), _EXAMPLE_MAX)
            sid = ids.get(a.shape_id, "")
            example = ""
            if a.is_candidate and a.type == "text" and text_val:
                example = _truncate(text_val, mc if mc > 0 else _EXAMPLE_MAX)
            shapes.append({
                "shape_id": a.shape_id, "name": a.name, "type": a.type,
                "bbox_pct": a.bbox_pct, "confidence": a.confidence,
                "is_candidate": a.is_candidate,
                "suggested_id": sid,
                "suggested_max_chars": mc, "suggested_max_lines": ml,
                "suggested_max_rows": mr, "suggested_max_cols": mcols,
                "font_pt": a.font_pt,
                "text": text_val,
                "suggested_example": example,
                "suggested_description": slot_description(sid),
            })
        cand = [a for a in assessments if a.is_candidate]
        text_blob = " ".join(
            _shape_text(shape_by_id[a.shape_id]) for a in cand if a.type == "text"
        )
        has_table = any(a.type == "table" for a in cand)
        num_text = sum(1 for a in cand if a.type == "text")
        has_subtitle = any(ids.get(a.shape_id) == "subtitle" for a in cand)
        slot_ids = [ids[a.shape_id] for a in cand if ids.get(a.shape_id)]
        kind = slide_kind(text_blob, has_table, i, num_text, has_subtitle)
        slides.append({
            "index": i, "width_emu": sw, "height_emu": sh, "shapes": shapes,
            "kind": kind, "suggested_name": kind,
            "suggested_description": slide_description(kind, slot_ids),
        })
    sigs = [slide_signature(s) for s in slides]
    counts = Counter(sigs)
    for s, sig in zip(slides, sigs):
        s["repeatable"] = counts[sig] >= 2
    return {"slides": slides}


def derive_ids(assessments: list[ShapeAssessment]) -> dict[int, str]:
    text = [a for a in assessments if a.type == "text"]
    by_area = sorted(text, key=lambda a: a.bbox_pct["w"] * a.bbox_pct["h"], reverse=True)
    ids: dict[int, str] = {}

    top_sorted = sorted(text, key=lambda a: a.bbox_pct["y"])
    title = top_sorted[0] if top_sorted else None
    if title is not None:
        ids[title.shape_id] = "title"
        below = [a for a in top_sorted[1:] if a.bbox_pct["y"] > title.bbox_pct["y"]]
        if below:
            ids[below[0].shape_id] = "subtitle"
    for a in by_area:
        if a.shape_id not in ids:
            ids[a.shape_id] = "body"
            break

    counters = {"text": 0, "table": 0, "image": 0}
    used = set(ids.values())
    for a in assessments:
        if a.shape_id in ids:
            continue
        base = "image" if a.type == "image" else a.type
        counters[base] = counters.get(base, 0) + 1
        cand = f"{base}_{counters[base]}"
        while cand in used:
            counters[base] += 1
            cand = f"{base}_{counters[base]}"
        ids[a.shape_id] = cand
        used.add(cand)
    return ids
