import copy

from pptx import Presentation

from .models import Template

# Namespace URI for Office Open XML relationships
_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Attribute name for the relationship id on <p:sldId> elements
_RID_ATTR = "{%s}id" % _REL_NS

# Reltype fragment that identifies a slide-layout relationship
_SLIDE_LAYOUT_RELTYPE_FRAGMENT = "slideLayout"

_EMBED_ATTR = "{%s}embed" % _REL_NS
_LINK_ATTR = "{%s}link" % _REL_NS
_ID_ATTR = "{%s}id" % _REL_NS
_REL_ATTRS = (_EMBED_ATTR, _LINK_ATTR, _ID_ATTR)


def _remap_rels(src_part, dest_part, element) -> dict:
    """Copy the relationships *element* references from src_part into dest_part,
    rewrite the r:embed/r:link/r:id ids on element in place, and return the
    old_rid -> new_rid map. Slide-layout rels are skipped (already wired via
    add_slide)."""
    used = set()
    for el in element.iter():
        for attr in _REL_ATTRS:
            val = el.get(attr)
            if val:
                used.add(val)

    old_to_new: dict[str, str] = {}
    for rid in used:
        if rid not in src_part.rels:
            continue
        rel = src_part.rels[rid]
        if _SLIDE_LAYOUT_RELTYPE_FRAGMENT in rel.reltype:
            continue
        if rel.is_external:
            old_to_new[rid] = dest_part.relate_to(rel._target, rel.reltype, is_external=True)
        else:
            old_to_new[rid] = dest_part.relate_to(rel._target, rel.reltype)

    if old_to_new:
        for el in element.iter():
            for attr in _REL_ATTRS:
                val = el.get(attr)
                if val and val in old_to_new:
                    el.set(attr, old_to_new[val])
    return old_to_new


def drop_base_slides(prs, count: int) -> None:
    """Remove the first *count* slides from the package: drop each <p:sldId>
    (ref-count -> 0) then drop_rel so the serialiser omits the orphaned Part."""
    xml_slides = prs.slides._sldIdLst
    originals = list(xml_slides)[:count]
    rids = [sid.get(_RID_ATTR) for sid in originals]
    for sid, rid in zip(originals, rids):
        xml_slides.remove(sid)
        if rid is not None:
            prs.part.drop_rel(rid)


def find_shape(slide, shape_id: int):
    """Return the shape on *slide* whose shape_id matches, or raise KeyError."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise KeyError(f"shape_id {shape_id} not found on slide")


def _duplicate_slide(prs: Presentation, src_index: int):
    """Append a copy of slide *src_index* (within the same package) and return it.

    Deep-copies every shape element from the source slide's spTree into the new
    slide.  Relationship ids in the *source* slide part may differ from those in
    the *destination* part, so this function:

    1. Copies every relationship from the source part into the dest part via
       ``Part.relate_to``, EXCEPT the slide-layout relationship (which was already
       wired when ``add_slide(layout)`` was called and must not be duplicated).
    2. Builds a complete ``old_rid -> new_rid`` map from those copies.
    3. Rewrites every ``r:embed``, ``r:link``, and ``r:id`` attribute in the
       copied XML tree using that map, so charts, hyperlinks, and images all
       resolve correctly in the destination part.
    """
    source = prs.slides[src_index]
    layout = source.slide_layout
    dest = prs.slides.add_slide(layout)

    # Remove any placeholder shapes that add_slide inserted from the layout.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # Deep-copy each source shape element into the dest spTree.
    for shp in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    # Copy + remap the relationships those shapes reference into the dest part.
    _remap_rels(source.part, dest.part, dest.shapes._spTree)

    return dest


def assemble(order: list[int], template: Template) -> Presentation:
    """Return a Presentation whose slides are copies of the template's base slides
    at the given *source_slide_index* values (in *order*), with the original base
    slides removed.

    All shape types (text, tables, pictures) survive a save -> reopen round-trip.

    Original slide Parts (and any media they exclusively own) are fully removed
    from the package by:
      1. Capturing each original slide's rId from its <p:sldId> element.
      2. Removing the <p:sldId> element from the presentation XML (drops the
         XML reference so _rel_ref_count falls to 0).
      3. Calling prs.part.drop_rel(rId) to remove the relationship entry and
         allow the serialiser to omit the orphaned Part from the saved zip.

    Because the duplicated slides hold their own independent relationships to
    shared image/media parts, dropping the originals does not break the copies.
    """
    prs = Presentation(template.pptx_path)
    original_count = len(prs.slides)

    for src_index in order:
        _duplicate_slide(prs, src_index)

    drop_base_slides(prs, original_count)
    return prs
