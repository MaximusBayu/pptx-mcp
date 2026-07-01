import copy
import io
import re

from pptx import Presentation
from pptx.oxml.ns import qn

from .assembler import drop_base_slides, find_shape, _remap_rels
from .catalog import component_type, get_catalog, _hex_or_none
from .filler import fill_shape
from .guardrails import check_layout
from .models import Constraints, SlotError, Template

_CID_RE = re.compile(r"^\d+:\d+$")
# catalog type -> predicate the content value must satisfy
_CONTENT_OK = {
    "text": lambda v: isinstance(v, str)
    or (isinstance(v, list) and all(isinstance(x, str) for x in v)),
    "table": lambda v: isinstance(v, list) and all(isinstance(r, list) for r in v),
    "image": lambda v: bool(v) and isinstance(v, (str, bytes)),
}


class ComposeRejected(Exception):
    def __init__(self, errors: list[SlotError]):
        self.errors = errors
        super().__init__(f"{len(errors)} composition error(s)")


def _num(v) -> bool:
    return isinstance(v, (int, float)) and not isinstance(v, bool)


def _bbox_ok(b) -> bool:
    if not isinstance(b, dict) or not all(k in b for k in ("x", "y", "w", "h")):
        return False
    if not all(_num(b[k]) for k in ("x", "y", "w", "h")):
        return False
    if not (0 <= b["x"] <= 100 and 0 <= b["y"] <= 100):
        return False
    return 0 < b["w"] <= 100 and 0 < b["h"] <= 100


def validate_composition(composition_spec: dict, template: Template) -> list[SlotError]:
    cat = get_catalog(template)
    by_id = {c["component_id"]: c for c in cat["components"]}
    n_slides = len(Presentation(template.pptx_path).slides)

    errors: list[SlotError] = []
    for i, slide in enumerate(composition_spec.get("slides", [])):
        canvas = slide.get("canvas")
        if not isinstance(canvas, int) or isinstance(canvas, bool) or not (0 <= canvas < n_slides):
            errors.append(SlotError(i, None, "unknown_canvas",
                                    f"canvas {canvas!r}; slides 0..{n_slides - 1}"))
        for placement in slide.get("placements", []):
            cid = placement.get("component_id")
            if not isinstance(cid, str) or not _CID_RE.match(cid) or cid not in by_id:
                errors.append(SlotError(i, cid, "unknown_component",
                                        f"no component {cid!r} in template"))
                continue
            ctype = by_id[cid]["type"]
            if "content" in placement and placement["content"] is not None:
                check = _CONTENT_OK.get(ctype)
                if check is None or not check(placement["content"]):
                    errors.append(SlotError(i, cid, "wrong_type",
                                            f"content not valid for {ctype} component"))
            if "bbox_pct" in placement and not _bbox_ok(placement["bbox_pct"]):
                errors.append(SlotError(i, cid, "bad_bbox",
                                        "bbox_pct needs numeric x,y(0-100), w,h(0-100, >0)"))
    return errors


def _copy_background(src_slide, dest_slide) -> None:
    """Copy the canvas slide's slide-level <p:bg> (if any) so the output slide
    matches the canvas background. If absent, the layout/master bg shows through.
    """
    src_csld = src_slide._element.find(qn("p:cSld"))
    if src_csld is None:
        return
    bg = src_csld.find(qn("p:bg"))
    if bg is None:
        return
    dest_csld = dest_slide._element.find(qn("p:cSld"))
    dest_csld.insert(0, copy.deepcopy(bg))  # schema: bg precedes spTree


def _set_geometry(shape, bbox, sw, sh) -> None:
    shape.left = int(sw * bbox["x"] / 100.0)
    shape.top = int(sh * bbox["y"] / 100.0)
    shape.width = int(sw * bbox["w"] / 100.0)
    shape.height = int(sh * bbox["h"] / 100.0)


def _rect_pct(shape, sw, sh) -> dict:
    return {"x": 100.0 * (shape.left or 0) / sw, "y": 100.0 * (shape.top or 0) / sh,
            "w": 100.0 * (shape.width or 0) / sw, "h": 100.0 * (shape.height or 0) / sh}


def _text_color(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    paras = shape.text_frame.paragraphs
    runs = paras[0].runs if paras else []
    return _hex_or_none(runs[0].font.color) if runs else None


def _eff_bg(shape, canvas_bg):
    try:
        fill = shape.fill
        if fill.type is not None:
            c = _hex_or_none(fill.fore_color)
            if c:
                return c
    except (TypeError, AttributeError, ValueError):
        pass
    return canvas_bg


def _canvas_bg_hex(slide):
    csld = slide._element.find(qn("p:cSld"))
    if csld is None:
        return None
    bg = csld.find(qn("p:bg"))
    if bg is None:
        return None
    clr = bg.find(".//" + qn("a:srgbClr"))
    return clr.get("val") if clr is not None else None


def compose(composition_spec: dict, template: Template) -> tuple[bytes, list[dict]]:
    errors = validate_composition(composition_spec, template)
    if errors:
        raise ComposeRejected(errors)

    prs = Presentation(template.pptx_path)
    original_count = len(prs.slides)
    base_slides = list(prs.slides)[:original_count]
    sw, sh = prs.slide_width, prs.slide_height

    # slot constraints keyed by (source_slide_index, shape_id) for fill defaults
    slot_map = {(st.source_slide_index, s.shape_id): s
                for st in template.slide_types for s in st.slots}

    warnings: list[dict] = []
    for out_index, slide_spec in enumerate(composition_spec["slides"]):
        canvas = base_slides[slide_spec["canvas"]]
        dest = prs.slides.add_slide(canvas.slide_layout)
        _copy_background(canvas, dest)
        canvas_bg = _canvas_bg_hex(canvas)
        for shp in list(dest.shapes):
            shp._element.getparent().remove(shp._element)

        placed_shapes = []  # (component_id, shape)
        for placement in slide_spec.get("placements", []):
            src_idx, shape_id = (int(x) for x in placement["component_id"].split(":"))
            src_shape = find_shape(base_slides[src_idx], shape_id)
            dest.shapes._spTree.append(copy.deepcopy(src_shape._element))
            _remap_rels(base_slides[src_idx].part, dest.part, dest.shapes._spTree[-1])
            placed = dest.shapes[-1]

            if "bbox_pct" in placement:
                _set_geometry(placed, placement["bbox_pct"], sw, sh)

            content = placement.get("content")
            if content is not None:
                kind = component_type(placed)
                slot = slot_map.get((src_idx, shape_id))
                constraints = slot.constraints if slot is not None else Constraints()
                try:
                    for w in fill_shape(dest, placed, kind, content, constraints,
                                        slot_id=placement["component_id"],
                                        max_bottom_emu=sh):
                        w.slide_index = out_index
                        warnings.append(w.to_dict())
                except Exception as exc:
                    warnings.append(SlotError(out_index, placement["component_id"],
                                              "fill_failed", str(exc)).to_dict())
                    continue
            placed_shapes.append((placement["component_id"], placed))

        gl = [{"component_id": cid, "rect": _rect_pct(shp, sw, sh),
               "text_color": _text_color(shp), "eff_bg": _eff_bg(shp, canvas_bg)}
              for cid, shp in placed_shapes]
        gwarnings, clamps = check_layout(gl)
        for w in gwarnings:
            w["slide_index"] = out_index
            warnings.append(w)
        for cid, shp in placed_shapes:
            if cid in clamps:
                _set_geometry(shp, clamps[cid], sw, sh)

    drop_base_slides(prs, original_count)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), warnings


def compose_dry_run(composition_spec: dict, template: Template) -> dict:
    """Validate + compose, discard the bytes; return errors and warnings."""
    try:
        _bytes, warnings = compose(composition_spec, template)
    except ComposeRejected as e:
        return {"errors": [err.to_dict() for err in e.errors], "warnings": []}
    return {"errors": [], "warnings": warnings}
