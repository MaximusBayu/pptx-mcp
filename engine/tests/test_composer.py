import io

from pptx import Presentation
from pptx_mcp.template import load_template
from pptx_mcp.catalog import get_catalog
from pptx_mcp.composer import compose


def _components(tpl):
    comps = get_catalog(tpl)["components"]
    return {
        "title": next(c for c in comps if c.get("slot_id") == "title"),
        "subtitle": next(c for c in comps if c.get("slot_id") == "subtitle"),
        "image": next(c for c in comps if c["type"] == "image"),
    }


def _reopen(data):
    return Presentation(io.BytesIO(data))


def test_single_placement_places_component(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "Hello World"}]}]}
    data, warnings = compose(spec, tpl)
    prs = _reopen(data)
    assert len(prs.slides) == 1
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert "Hello World" in texts


def test_cross_slide_image_survives_roundtrip(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["image"]["component_id"],
         "bbox_pct": {"x": 10, "y": 10, "w": 40, "h": 40}}]}]}
    data, _ = compose(spec, tpl)
    prs = _reopen(data)
    pics = [s for s in prs.slides[0].shapes if s.shape_type == 13]
    assert pics, "cloned picture missing after save/reopen"


def test_manifest_model_drops_unmentioned(sample_template_dir):
    # canvas 0 has TITLE + SUBTITLE; only place TITLE -> SUBTITLE must be gone.
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"]}]}]}
    data, _ = compose(spec, tpl)
    prs = _reopen(data)
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert not any("SUBTITLE" in t for t in texts)


def test_verbatim_clone_keeps_sample_text(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"]}]}]}  # no content
    data, _ = compose(spec, tpl)
    prs = _reopen(data)
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert any("TITLE" in t for t in texts)


def test_oversized_text_warns(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "First sentence. " * 30}]}]}
    _data, warnings = compose(spec, tpl)
    assert any(w["code"] == "text_truncated" for w in warnings)


def test_zorder_follows_placement_order(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "AAA"},
        {"component_id": c["subtitle"]["component_id"], "content": "BBB"}]}]}
    data, _ = compose(spec, tpl)
    prs = _reopen(data)
    last = prs.slides[0].shapes[-1]
    assert last.has_text_frame and last.text_frame.text == "BBB"


def test_bbox_repositions(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    src = Presentation(str(sample_template_dir / "base.pptx"))
    sw = src.slide_width
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "X",
         "bbox_pct": {"x": 50, "y": 50, "w": 25, "h": 10}}]}]}
    data, _ = compose(spec, tpl)
    prs = _reopen(data)
    shp = next(s for s in prs.slides[0].shapes if s.has_text_frame)
    assert abs(shp.left - int(sw * 0.5)) < int(sw * 0.01)


def test_list_content_renders_bullets_end_to_end(sample_template_dir):
    tpl = load_template(sample_template_dir)
    comps = get_catalog(tpl)["components"]
    body = next(c for c in comps if c.get("slot_id") == "body")
    spec = {"slides": [{"canvas": 1, "placements": [
        {"component_id": body["component_id"], "content": ["First", "Second", "Third"]}]}]}
    data, _ = compose(spec, tpl)
    prs = _reopen(data)
    body_shapes = [s for s in prs.slides[0].shapes
                   if s.has_text_frame and len(s.text_frame.paragraphs) == 3]
    assert body_shapes, "expected a 3-paragraph bullet box"


def test_fill_failure_becomes_warning_not_crash(sample_template_dir, monkeypatch):
    import pptx_mcp.composer as comp
    tpl = load_template(sample_template_dir)
    c = _components(tpl)

    def boom(*a, **k):
        raise RuntimeError("kaboom")

    monkeypatch.setattr(comp, "fill_shape", boom)
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "X"}]}]}
    data, warnings = compose(spec, tpl)
    assert data[:2] == b"PK"
    assert any(w["code"] == "fill_failed" for w in warnings)


def test_off_slide_placement_clamped(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    src = Presentation(str(sample_template_dir / "base.pptx"))
    sw = src.slide_width
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "X",
         "bbox_pct": {"x": 90, "y": 10, "w": 30, "h": 10}}]}]}
    data, warnings = compose(spec, tpl)
    prs = _reopen(data)
    shp = next(s for s in prs.slides[0].shapes if s.has_text_frame)
    assert shp.left + shp.width <= sw + 1  # clamped inside slide
    assert any(w["code"] == "clamped" for w in warnings)


def test_overlapping_placements_warn(sample_template_dir):
    tpl = load_template(sample_template_dir)
    c = _components(tpl)
    box = {"x": 10, "y": 10, "w": 50, "h": 50}
    spec = {"slides": [{"canvas": 0, "placements": [
        {"component_id": c["title"]["component_id"], "content": "A", "bbox_pct": box},
        {"component_id": c["subtitle"]["component_id"], "content": "B", "bbox_pct": box}]}]}
    _data, warnings = compose(spec, tpl)
    assert any(w["code"] == "overlap" for w in warnings)
