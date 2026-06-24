import io
import pytest
from pptx import Presentation
from pptx_mcp.template import load_template
from pptx_mcp.assembler import find_shape
from pptx_mcp.render import render, RenderRejected, dry_run


def _deck():
    return {"slides": [
        {"slide_type": "title", "slots": {"title": "Acme", "subtitle": "Q3"}},
        {"slide_type": "bullet", "slots": {"heading": "Plan", "body": "a\nb"}},
    ]}


def test_render_produces_valid_pptx(sample_template_dir):
    tpl = load_template(sample_template_dir)
    data, warnings = render(_deck(), tpl)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 2
    texts = [sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame]
    assert "Acme" in texts


def test_render_rejects_invalid(sample_template_dir):
    tpl = load_template(sample_template_dir)
    # Tables still reject; text overflow is now a non-fatal warning
    bad = {"slides": [{"slide_type": "table", "slots": {"data": [["a"]] * 10}}]}
    with pytest.raises(RenderRejected) as ei:
        render(bad, tpl)
    assert ei.value.errors[0].code == "table_overflow"


def test_render_text_overflow_returns_warnings(sample_template_dir):
    tpl = load_template(sample_template_dir)
    # Text over limit now produces a warning instead of raising
    long_text = "First sentence here. " * 20  # well over max_chars=40
    spec = {"slides": [{"slide_type": "title", "slots": {"title": long_text, "subtitle": ""}}]}
    data, warnings = render(spec, tpl)
    assert isinstance(data, bytes)
    assert any(w["code"] == "text_truncated" for w in warnings)


def test_render_clears_omitted_optional_text_slot(sample_template_dir):
    """Omitting an optional text slot must clear the template's sample text."""
    tpl = load_template(sample_template_dir)
    # title slide: "title" is required, "subtitle" is optional (required=False).
    # Provide only the required slot; omit subtitle entirely.
    st = tpl.slide_type("title")
    opt = next(s for s in st.slots if s.type == "text" and not s.required)
    deck = {"slides": [{"slide_type": "title", "slots": {"title": "Hello"}}]}
    pptx_bytes, warnings = render(deck, tpl)
    prs = Presentation(io.BytesIO(pptx_bytes))
    shp = find_shape(prs.slides[0], opt.shape_id)
    assert shp.text_frame.text == "", (
        f"Expected subtitle shape to be cleared but got: {shp.text_frame.text!r}"
    )


def test_dry_run_valid_deck_returns_no_errors(sample_template_dir):
    tpl = load_template(sample_template_dir)
    # title slide: "title" is required, "subtitle" is optional.
    # Providing both required slots makes this a genuinely valid deck.
    deck = {"slides": [{"slide_type": "title", "slots": {"title": "Acme", "subtitle": "Q3"}}]}
    result = dry_run(deck, tpl)
    assert result["errors"] == []
    assert isinstance(result["warnings"], list)


def test_dry_run_invalid_deck_returns_errors(sample_template_dir):
    tpl = load_template(sample_template_dir)
    deck = {"slides": [{"slide_type": "does_not_exist", "slots": {}}]}
    result = dry_run(deck, tpl)
    # An unknown slide_type yields validation errors, no warnings, and never raises.
    assert len(result["errors"]) >= 1
    assert result["warnings"] == []
