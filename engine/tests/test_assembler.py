import io

from pptx import Presentation
from pptx_mcp.template import load_template
from pptx_mcp.assembler import assemble, find_shape

# Namespace URI for Office Open XML relationships
_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def test_assemble_order_and_count(sample_template_dir):
    tpl = load_template(sample_template_dir)
    # build: title(0), bullet(1), title(0) again -> 3 slides
    prs = assemble([0, 1, 0], tpl)
    assert len(prs.slides) == 3


def test_assembled_slides_keep_shapes(sample_template_dir):
    """Round-trip test: assemble, save to BytesIO, reopen, assert shapes survive.

    This test proves that the rId rewrite in _duplicate_slide actually survives
    serialisation: it opens the *reopened* Presentation, not the in-memory one.
    It additionally confirms the picture's blip relationship resolves to a real
    image part with non-empty blob data.
    """
    tpl = load_template(sample_template_dir)
    prs = assemble([2, 3], tpl)  # table slide (idx 0), image slide (idx 1)

    # Save to BytesIO and reopen — this is the round-trip that proves the spike.
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    reopened = Presentation(buf)

    table_slide = reopened.slides[0]
    image_slide = reopened.slides[1]

    # Table slide must still have a table.
    assert any(s.has_table for s in table_slide.shapes), (
        "table slide lost its table after save/reopen"
    )

    # Image slide must still have a picture shape (shape_type 13 = PICTURE).
    pic_shapes = [s for s in image_slide.shapes if s.shape_type == 13]
    assert pic_shapes, "image slide lost its picture shape after save/reopen"

    # Verify the blip relationship actually resolves: find the r:embed rId in the
    # picture XML and confirm the corresponding relationship target has blob data.
    embed_attr = "{%s}embed" % _REL_NS
    slide_part = image_slide.part
    embed_rids = set()
    for elem in image_slide.shapes._spTree.iter():
        val = elem.get(embed_attr)
        if val:
            embed_rids.add(val)

    assert embed_rids, "no r:embed attribute found in image slide XML"
    for embed_rid in embed_rids:
        assert embed_rid in slide_part.rels, (
            f"r:embed rId '{embed_rid}' not present in image slide relationships"
        )
        rel = slide_part.rels[embed_rid]
        assert not rel.is_external, f"r:embed rId '{embed_rid}' is unexpectedly external"
        target_part = rel._target
        assert hasattr(target_part, "blob"), (
            f"relationship target for rId '{embed_rid}' has no blob"
        )
        assert len(target_part.blob) > 0, (
            f"image blob for rId '{embed_rid}' is empty after save/reopen"
        )


def test_find_shape_by_id(sample_template_dir):
    tpl = load_template(sample_template_dir)
    prs = assemble([0], tpl)
    title_slot = tpl.slide_type("title").slot("title")
    shp = find_shape(prs.slides[0], title_slot.shape_id)
    assert shp.shape_id == title_slot.shape_id
    # Assert the shape's text equals "TITLE" so the lookup is not tautological.
    assert shp.text_frame.text == "TITLE", (
        f"expected shape text 'TITLE', got {shp.text_frame.text!r}"
    )


def test_remap_rels_copies_referenced_image_rel(sample_template_dir):
    """A single picture element deep-copied to a fresh slide has its blip rel
    remapped into the dest part."""
    import copy
    from pptx import Presentation
    from pptx_mcp.assembler import _remap_rels

    prs = Presentation(str(sample_template_dir / "base.pptx"))
    src_slide = prs.slides[3]  # image slide
    pic = next(s for s in src_slide.shapes if s.shape_type == 13)

    dest = prs.slides.add_slide(prs.slide_layouts[6])
    el = copy.deepcopy(pic._element)
    dest.shapes._spTree.append(el)

    mapping = _remap_rels(src_slide.part, dest.part, el)
    assert mapping  # at least the blip rel was remapped

    embed_attr = "{%s}embed" % "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    new_rids = {e.get(embed_attr) for e in el.iter() if e.get(embed_attr)}
    assert new_rids
    for rid in new_rids:
        assert rid in dest.part.rels  # resolves in the dest part
