import io

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from pptx_mcp.filler import fill_slot
from pptx_mcp.models import Constraints, Slot


def _deck_with_styled_textbox():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(3000000), Emu(800000))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PLACEHOLDER"
    run.font.bold = True
    run.font.italic = True
    run.font.name = "Georgia"
    run.font.size = Pt(18)
    return prs, slide, tb


def _slot(shape_id, constraints=None):
    return Slot(id="title", name="Title", type="text", shape_id=shape_id,
                constraints=constraints or Constraints())


def test_fill_preserves_alignment_font_bold_italic():
    prs, slide, tb = _deck_with_styled_textbox()
    fill_slot(slide, _slot(tb.shape_id), "New Heading")
    p = tb.text_frame.paragraphs[0]
    run = p.runs[0]
    assert run.text == "New Heading"
    assert p.alignment == PP_ALIGN.CENTER
    assert run.font.bold is True
    assert run.font.italic is True
    assert run.font.name == "Georgia"
    assert run.font.size == Pt(18)
    # No stray leftover runs/paragraphs.
    assert len(tb.text_frame.paragraphs) == 1
    assert len(p.runs) == 1


def test_shrink_scales_from_original_size_not_hardcoded():
    # Force a shrink with a tight max_chars; original is 18pt, must shrink below it.
    prs, slide, tb = _deck_with_styled_textbox()
    fill_slot(slide, _slot(tb.shape_id, Constraints(max_chars=5)),
              "This is a long heading that overflows")
    run = tb.text_frame.paragraphs[0].runs[0]
    assert run.font.size is not None
    assert run.font.size.pt < 18  # shrunk relative to the template's size
    # Styling still preserved through the shrink path.
    assert run.font.bold is True
    assert tb.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER


def test_empty_box_falls_back_without_crash():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(2000000), Emu(600000))
    # no runs added
    fill_slot(slide, _slot(tb.shape_id), "Fallback text")
    assert tb.text_frame.text == "Fallback text"
