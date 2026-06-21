import io
from pptx import Presentation
from pptx_mcp.move import move_shape, move_shapes


def _first_shape(sample_template_dir):
    prs = Presentation(str(sample_template_dir / "base.pptx"))
    return prs.slides[0].shapes[0].shape_id


def test_move_shape_repositions(sample_template_dir):
    pptx = (sample_template_dir / "base.pptx").read_bytes()
    sid = _first_shape(sample_template_dir)
    out = move_shape(pptx, sid, {"x": 10, "y": 20, "w": 50, "h": 25})
    prs = Presentation(io.BytesIO(out))
    sw, sh = prs.slide_width, prs.slide_height
    moved = next(s for s in prs.slides[0].shapes if s.shape_id == sid)
    assert abs(moved.left - sw * 0.10) < 5000
    assert abs(moved.top - sh * 0.20) < 5000


def _two_slide_deck() -> bytes:
    """Two blank slides, each with one textbox. On a blank layout the first
    added shape gets shape_id=2 on both slides -> the cross-slide collision we
    must handle."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(2):
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(0, 0, 914400, 914400)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_move_shapes_targets_correct_slide():
    data = _two_slide_deck()
    prs = Presentation(io.BytesIO(data))
    sid1 = prs.slides[1].shapes[0].shape_id
    out = move_shapes(data, [
        {"slide_index": 1, "shape_id": sid1,
         "bbox_pct": {"x": 50, "y": 50, "w": 20, "h": 10}},
    ])
    prs2 = Presentation(io.BytesIO(out))
    sw, sh = prs2.slide_width, prs2.slide_height
    s1 = prs2.slides[1].shapes[0]
    assert abs(s1.left - sw * 0.50) < 5000
    assert abs(s1.top - sh * 0.50) < 5000
    # Slide 0's same-id shape is untouched (still at origin).
    assert prs2.slides[0].shapes[0].left == 0


def test_move_shapes_applies_batch():
    data = _two_slide_deck()
    prs = Presentation(io.BytesIO(data))
    sid0 = prs.slides[0].shapes[0].shape_id
    sid1 = prs.slides[1].shapes[0].shape_id
    out = move_shapes(data, [
        {"slide_index": 0, "shape_id": sid0, "bbox_pct": {"x": 10, "y": 10, "w": 30, "h": 15}},
        {"slide_index": 1, "shape_id": sid1, "bbox_pct": {"x": 60, "y": 20, "w": 30, "h": 15}},
    ])
    prs2 = Presentation(io.BytesIO(out))
    sw = prs2.slide_width
    assert abs(prs2.slides[0].shapes[0].left - sw * 0.10) < 5000
    assert abs(prs2.slides[1].shapes[0].left - sw * 0.60) < 5000


def test_move_shapes_unknown_raises():
    import pytest
    data = _two_slide_deck()
    with pytest.raises(KeyError):
        move_shapes(data, [{"slide_index": 5, "shape_id": 2,
                            "bbox_pct": {"x": 1, "y": 1, "w": 1, "h": 1}}])


def test_move_shape_slide_index_scopes_search():
    data = _two_slide_deck()
    prs = Presentation(io.BytesIO(data))
    sid = prs.slides[1].shapes[0].shape_id
    out = move_shape(data, sid, {"x": 40, "y": 0, "w": 20, "h": 10}, slide_index=1)
    prs2 = Presentation(io.BytesIO(out))
    sw = prs2.slide_width
    assert abs(prs2.slides[1].shapes[0].left - sw * 0.40) < 5000
    assert prs2.slides[0].shapes[0].left == 0
