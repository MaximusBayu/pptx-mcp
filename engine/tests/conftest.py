import io
import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


def _add_textbox(slide, left, top, width, height, text):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    return box


def _build_pptx(path: Path) -> dict:
    """Build a 4-slide base.pptx; return shape ids per slide."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    ids = {}

    # slide 0: title
    s0 = prs.slides.add_slide(blank)
    title = _add_textbox(s0, 1, 1, 8, 1.5, "TITLE")
    subtitle = _add_textbox(s0, 1, 2.6, 8, 1, "SUBTITLE")
    ids["title"] = {"title": title.shape_id, "subtitle": subtitle.shape_id}

    # slide 1: bullet
    s1 = prs.slides.add_slide(blank)
    heading = _add_textbox(s1, 1, 0.5, 8, 1, "HEADING")
    body = _add_textbox(s1, 1, 1.6, 8, 4, "BODY")
    ids["bullet"] = {"heading": heading.shape_id, "body": body.shape_id}

    # slide 2: table (2x2 placeholder)
    s2 = prs.slides.add_slide(blank)
    gf = s2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(8), Inches(3))
    ids["table"] = {"data": gf.shape_id}

    # slide 3: image (a starter picture)
    s3 = prs.slides.add_slide(blank)
    img_bytes = _tiny_png()
    pic = s3.shapes.add_picture(io.BytesIO(img_bytes), Inches(1), Inches(1), Inches(4), Inches(3))
    ids["image"] = {"photo": pic.shape_id}

    prs.save(str(path))
    return ids


def _tiny_png() -> bytes:
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", (10, 10), (200, 50, 50)).save(buf, format="PNG")
    return buf.getvalue()


def _build_manifest(ids: dict) -> dict:
    return {
        "template": {"id": "sample", "name": "Sample", "description": "Test template"},
        "slide_types": [
            {
                "id": "title", "name": "Title Slide", "description": "Opening",
                "source_slide_index": 0,
                "slots": [
                    {"id": "title", "name": "Title", "type": "text",
                     "target": {"shape_id": ids["title"]["title"]},
                     "required": True, "default": None,
                     "constraints": {"max_chars": 40, "max_lines": 2, "shrink_floor_pt": 16}},
                    {"id": "subtitle", "name": "Subtitle", "type": "text",
                     "target": {"shape_id": ids["title"]["subtitle"]},
                     "required": False, "default": "",
                     "constraints": {"max_chars": 60, "max_lines": 2, "shrink_floor_pt": 12}},
                ],
            },
            {
                "id": "bullet", "name": "Bullet Slide", "description": "Bulleted body",
                "source_slide_index": 1,
                "slots": [
                    {"id": "heading", "name": "Heading", "type": "text",
                     "target": {"shape_id": ids["bullet"]["heading"]},
                     "required": True, "default": None,
                     "constraints": {"max_chars": 50, "max_lines": 1, "shrink_floor_pt": 16}},
                    {"id": "body", "name": "Body", "type": "text",
                     "target": {"shape_id": ids["bullet"]["body"]},
                     "required": True, "default": None,
                     "constraints": {"max_chars": 300, "max_lines": 8, "shrink_floor_pt": 12}},
                ],
            },
            {
                "id": "table", "name": "Table Slide", "description": "Data table",
                "source_slide_index": 2,
                "slots": [
                    {"id": "data", "name": "Data", "type": "table",
                     "target": {"shape_id": ids["table"]["data"]},
                     "required": True, "default": None,
                     "constraints": {"max_rows": 5, "max_cols": 4}},
                ],
            },
            {
                "id": "image", "name": "Image Slide", "description": "Photo",
                "source_slide_index": 3,
                "slots": [
                    {"id": "photo", "name": "Photo", "type": "image",
                     "target": {"shape_id": ids["image"]["photo"]},
                     "required": True, "default": None,
                     "constraints": {"fit": "cover"}},
                ],
            },
        ],
    }


@pytest.fixture
def sample_template_dir(tmp_path) -> Path:
    ids = _build_pptx(tmp_path / "base.pptx")
    manifest = _build_manifest(ids)
    (tmp_path / "manifest.json").write_text(json.dumps(manifest, indent=2))
    return tmp_path


@pytest.fixture
def sample_manifest(sample_template_dir) -> dict:
    return json.loads((sample_template_dir / "manifest.json").read_text())


@pytest.fixture
def tiny_png_bytes() -> bytes:
    return _tiny_png()


@pytest.fixture
def base_template(tmp_path):
    """1-slide template with one text slot for overflow tests."""
    from pptx_mcp.bytesio import load_from_bytes

    prs = Presentation()
    blank = prs.slide_layouts[6]
    s0 = prs.slides.add_slide(blank)
    box = _add_textbox(s0, 1, 1, 8, 1.5, "CONTENT")
    shape_id = box.shape_id

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    manifest = {
        "template": {"id": "base", "name": "Base", "description": "Overflow test template"},
        "slide_types": [
            {
                "id": "content", "name": "Content Slide", "description": "Single text slot",
                "source_slide_index": 0,
                "slots": [
                    {"id": "body", "name": "Body", "type": "text",
                     "target": {"shape_id": shape_id},
                     "required": True, "default": None,
                     "constraints": {"max_chars": 100, "max_lines": 5, "shrink_floor_pt": 12}},
                ],
            },
        ],
    }
    return load_from_bytes(pptx_bytes, manifest)


@pytest.fixture
def labeled_deck(tmp_path):
    """A 1-slide deck with known slots + decoration. Returns (path, labels).
    labels maps shape_id -> True (content slot) / False (decoration)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    labels = {}

    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.5))
    title.text_frame.text = "Quarterly Business Review"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    labels[title.shape_id] = True

    body = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    body.text_frame.text = "Lorem ipsum dolor sit amet, consectetur adipiscing."
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    labels[body.shape_id] = True

    tiny = slide.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(0.2), Inches(0.2))
    labels[tiny.shape_id] = False

    line = slide.shapes.add_connector(2, Inches(1), Inches(6.9), Inches(9), Inches(6.9))
    labels[line.shape_id] = False

    p = tmp_path / "labeled.pptx"
    prs.save(str(p))
    return str(p), labels
