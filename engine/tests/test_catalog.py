from pptx_mcp.template import load_template
from pptx_mcp.catalog import get_catalog, _component_dict, _shape_style


def test_catalog_lists_all_components(sample_template_dir):
    tpl = load_template(sample_template_dir)
    cat = get_catalog(tpl)
    assert cat["id"] == "sample"
    comps = cat["components"]
    assert len(comps) >= 4  # at least one shape per sample slide
    # the title-slide title shape is a fillable slot
    title = next(c for c in comps if c["slot_id"] == "title")
    assert title["fillable"] is True
    assert title["component_id"] == f"{title['source_slide']}:" + title["component_id"].split(":")[1]
    g = title["geometry"]
    assert set(g["bbox_pct"]) == {"x", "y", "w", "h"}
    assert g["width_emu"] > 0 and g["height_emu"] > 0


def test_catalog_marks_types(sample_template_dir):
    tpl = load_template(sample_template_dir)
    comps = get_catalog(tpl)["components"]
    assert any(c["type"] == "table" for c in comps)   # slide 2 table
    assert any(c["type"] == "image" for c in comps)   # slide 3 picture


def test_component_dict_fillable_flag():
    class _Shp:
        shape_id = 7; name = "Box"; left = 0; top = 0; width = 100; height = 50
        shape_type = None; has_table = False; has_text_frame = False
    decor = _component_dict(_Shp(), 1, 1000, 1000, None)
    assert decor["fillable"] is False and decor["slot_id"] is None
    assert decor["component_id"] == "1:7"
    slot = _component_dict(_Shp(), 0, 1000, 1000, "title")
    assert slot["fillable"] is True and slot["slot_id"] == "title"


def test_shape_style_best_effort_no_crash():
    class _NoText:
        has_text_frame = False
        @property
        def fill(self):  # a fill whose fore_color has no rgb
            raise ValueError("no fill")
    style = _shape_style(_NoText())
    assert style == {"font_name": None, "font_pt": None,
                     "font_color": None, "fill_color": None}


def test_multiline_flag_and_hint(sample_template_dir):
    from pptx_mcp.template import load_template
    from pptx_mcp.catalog import get_catalog
    tpl = load_template(sample_template_dir)
    comps = get_catalog(tpl)["components"]
    for c in comps:
        assert "multiline" in c and isinstance(c["multiline"], bool)
        assert "hint" in c and isinstance(c["hint"], str) and c["hint"]
    img = next(c for c in comps if c["type"] == "image")
    assert "URL" in img["hint"] or "base64" in img["hint"]
    table = next(c for c in comps if c["type"] == "table")
    assert "list[list]" in table["hint"]


def test_multiline_true_for_bulleted_text():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    from pptx_mcp.catalog import _is_multiline
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "one"
    r.font.size = Pt(18)
    pPr = tb.text_frame.paragraphs[0]._p.get_or_add_pPr()
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "•"}))
    assert _is_multiline(tb) is True


def test_multiline_false_for_single_line():
    from pptx import Presentation
    from pptx.util import Inches
    from pptx_mcp.catalog import _is_multiline
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "just one line"
    assert _is_multiline(tb) is False
