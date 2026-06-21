import io

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from pptx_mcp.filler import _fill_image


def _img(w, h) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (w, h), (10, 120, 90)).save(buf, format="PNG")
    return buf.getvalue()


def _deck_with_box(box_w_emu, box_h_emu):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(box_w_emu), Emu(box_h_emu))
    return prs, slide, shape


def _only_picture(slide):
    return next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)


def test_contain_wide_image_in_tall_box_centers_vertically():
    # box 1x2 (tall), image 4x1 (wide) -> fills width, shrinks height, centered top.
    prs, slide, shape = _deck_with_box(914400, 1828800)
    box_left, box_top, box_w, box_h = shape.left, shape.top, shape.width, shape.height
    _fill_image(slide, shape, _img(400, 100), "contain")
    pic = _only_picture(slide)
    assert pic.width == box_w
    assert pic.height < box_h
    assert pic.left == box_left
    assert pic.top > box_top  # centered down


def test_contain_tall_image_in_wide_box_centers_horizontally():
    prs, slide, shape = _deck_with_box(1828800, 914400)
    box_left, box_top, box_w, box_h = shape.left, shape.top, shape.width, shape.height
    _fill_image(slide, shape, _img(100, 400), "contain")
    pic = _only_picture(slide)
    assert pic.height == box_h
    assert pic.width < box_w
    assert pic.top == box_top
    assert pic.left > box_left


def test_square_image_in_square_box_fills():
    prs, slide, shape = _deck_with_box(914400, 914400)
    box_w, box_h = shape.width, shape.height
    _fill_image(slide, shape, _img(200, 200), "contain")
    pic = _only_picture(slide)
    assert abs(pic.width - box_w) < 2
    assert abs(pic.height - box_h) < 2


def test_unreadable_image_falls_back_to_box_rect():
    prs, slide, shape = _deck_with_box(914400, 914400)
    box_left, box_top, box_w, box_h = shape.left, shape.top, shape.width, shape.height
    _fill_image(slide, shape, b"not-an-image", "contain")
    pic = _only_picture(slide)
    assert (pic.left, pic.top, pic.width, pic.height) == (box_left, box_top, box_w, box_h)
