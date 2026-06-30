from pptx.util import Pt
from pptx_mcp.template import load_template
from pptx_mcp.assembler import assemble, find_shape
from pptx_mcp.filler import fill_slot


def test_fill_text(sample_template_dir):
    tpl = load_template(sample_template_dir)
    prs = assemble([0], tpl)
    slot = tpl.slide_type("title").slot("title")
    fill_slot(prs.slides[0], slot, "Hello World")
    shp = find_shape(prs.slides[0], slot.shape_id)
    assert shp.text_frame.text == "Hello World"


def test_fill_text_sets_word_wrap_and_preserves_styling(sample_template_dir):
    tpl = load_template(sample_template_dir)
    prs = assemble([0], tpl)
    slot = tpl.slide_type("title").slot("title")
    shp = find_shape(prs.slides[0], slot.shape_id)
    before_name = shp.text_frame.paragraphs[0].runs[0].font.name

    fill_slot(prs.slides[0], slot, "A reasonably long heading that should wrap and be fit to its box")

    tf = shp.text_frame
    assert tf.word_wrap is True
    p0 = tf.paragraphs[0]
    # Font name (family) preserved on the surviving run.
    assert p0.runs[0].font.name == before_name
    # Line spacing is resolved to a numeric multiple and never below the floor.
    from pptx_mcp.textfit import LINE_SPACING_FLOOR
    assert isinstance(p0.line_spacing, float)
    assert p0.line_spacing >= LINE_SPACING_FLOOR


def test_fill_text_no_warning_when_it_fits(sample_template_dir):
    tpl = load_template(sample_template_dir)
    prs = assemble([0], tpl)
    slot = tpl.slide_type("title").slot("title")
    warnings = fill_slot(prs.slides[0], slot, "Short")
    assert not any(w.code == "text_truncated" for w in warnings)


def test_fill_table(sample_template_dir):
    tpl = load_template(sample_template_dir)
    prs = assemble([2], tpl)
    slot = tpl.slide_type("table").slot("data")
    fill_slot(prs.slides[0], slot, [["A", "B"], ["1", "2"]])
    shp = find_shape(prs.slides[0], slot.shape_id)
    assert shp.table.cell(0, 0).text == "A"
    assert shp.table.cell(1, 1).text == "2"


def test_fill_image(sample_template_dir, tiny_png_bytes):
    tpl = load_template(sample_template_dir)
    prs = assemble([3], tpl)
    slot = tpl.slide_type("image").slot("photo")
    fill_slot(prs.slides[0], slot, tiny_png_bytes)  # should not raise
    shp = find_shape(prs.slides[0], slot.shape_id)
    assert shp.shape_type == 13  # still a picture


def test_fill_text_overflow_cuts_and_reports(base_template):
    from pptx import Presentation
    from pptx_mcp.filler import fill_slot
    tpl = base_template
    slot = tpl.slide_types[0].slots[0]
    slot.constraints.max_chars = 10
    prs = Presentation(tpl.pptx_path)
    warnings = fill_slot(prs.slides[0], slot, "First short. Second sentence dropped.")
    assert any(w.code == "text_truncated" for w in warnings)


import math

from pptx import Presentation
from pptx.util import Emu, Pt

from pptx_mcp.filler import (
    fill_slot, _fill_table, _any_cell_overflows, _fit_cell,
)
from pptx_mcp.models import Constraints, Slot


def _deck_with_table(rows, cols, col_w=1_000_000, row_h=400_000):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gf = slide.shapes.add_table(rows, cols, Emu(500_000), Emu(500_000),
                                Emu(col_w * cols), Emu(row_h * rows))
    table = gf.table
    for c in range(cols):
        table.columns[c].width = Emu(col_w)
    for r in range(rows):
        table.rows[r].height = Emu(row_h)
    return prs, slide, gf, table


def _table_slot(shape_id):
    return Slot(id="table_1", name="Table", type="table", shape_id=shape_id,
                constraints=Constraints())


def test_short_cells_no_resize_no_warning():
    prs, slide, gf, table = _deck_with_table(2, 2)
    before_w = [table.columns[c].width for c in range(2)]
    warnings = _fill_table(gf, [["a", "b"], ["c", "d"]])
    after_w = [table.columns[c].width for c in range(2)]
    assert after_w == before_w          # no overflow -> no resize
    assert warnings == []
    assert table.cell(0, 0).text == "a"


def test_footprint_constant_after_resize():
    prs, slide, gf, table = _deck_with_table(2, 2)
    total_w_before = sum(table.columns[c].width for c in range(2))
    total_h_before = sum(table.rows[r].height for r in range(2))
    long = "x" * 4000
    _fill_table(gf, [[long, "b"], ["c", "d"]])
    total_w_after = sum(table.columns[c].width for c in range(2))
    total_h_after = sum(table.rows[r].height for r in range(2))
    assert total_w_after == total_w_before
    assert total_h_after == total_h_before


def test_overflowing_column_widened_over_slack():
    prs, slide, gf, table = _deck_with_table(1, 2)
    _fill_table(gf, [["x" * 4000, "b"]])
    assert table.columns[0].width > table.columns[1].width


def test_long_cell_truncated_with_warning():
    prs, slide, gf, table = _deck_with_table(1, 1, col_w=300_000, row_h=200_000)
    warnings = _fill_table(gf, [["y" * 5000]])
    assert any(w.code == "text_truncated" for w in warnings)


def test_fill_slot_propagates_table_warnings():
    prs, slide, gf, table = _deck_with_table(1, 1, col_w=300_000, row_h=200_000)
    warnings = fill_slot(slide, _table_slot(gf.shape_id), [["y" * 5000]])
    assert any(w.code == "text_truncated" for w in warnings)


def test_fit_cell_short_value_unchanged_font():
    prs, slide, gf, table = _deck_with_table(1, 1)
    cell = table.cell(0, 0)
    cell.text_frame.paragraphs[0].add_run().text = "old"
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    warnings = _fit_cell(cell, "hi", Emu(1_000_000), Emu(400_000), "cell[0,0]")
    assert warnings == []
    assert cell.text == "hi"
    assert cell.text_frame.paragraphs[0].runs[0].font.size == Pt(18)


def test_fit_cell_fitting_theme_cell_keeps_inherited_font():
    prs, slide, gf, table = _deck_with_table(1, 1)
    cell = table.cell(0, 0)
    run = cell.text_frame.paragraphs[0].add_run()
    run.text = "old"  # no explicit font.size -> inherits from theme
    _fit_cell(cell, "hi", Emu(1_000_000), Emu(400_000), "cell[0,0]")
    assert cell.text == "hi"
    assert cell.text_frame.paragraphs[0].runs[0].font.size is None  # not pinned


from pptx import Presentation
from pptx.util import Emu
from pptx_mcp.filler import clear_slot
from pptx_mcp.models import Constraints, Slot


def _slot(shape_id, type_):
    return Slot(id=f"{type_}_1", name=type_, type=type_, shape_id=shape_id,
                constraints=Constraints())


def test_clear_slot_text_empties_text_frame():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(2000000), Emu(600000))
    tb.text_frame.paragraphs[0].add_run().text = "TEMPLATE SAMPLE"
    clear_slot(slide, _slot(tb.shape_id, "text"))
    assert tb.text_frame.text == ""


def test_clear_slot_image_removes_shape():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(914400), Emu(914400), Emu(1000000), Emu(1000000))
    sid = tb.shape_id
    clear_slot(slide, _slot(sid, "image"))
    assert all(shp.shape_id != sid for shp in slide.shapes)


def test_clear_slot_table_blanks_all_cells():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gf = slide.shapes.add_table(2, 2, Emu(500000), Emu(500000), Emu(4000000), Emu(1000000))
    table = gf.table
    for r in range(2):
        for c in range(2):
            table.cell(r, c).text = "sample"
    clear_slot(slide, _slot(gf.shape_id, "table"))
    assert all(table.cell(r, c).text == "" for r in range(2) for c in range(2))


def test_clear_slot_missing_shape_is_noop():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # shape_id 99999 does not exist -> find_shape raises KeyError -> no-op, no raise
    clear_slot(slide, _slot(99999, "text"))


def test_fill_table_partial_blanks_surplus_template_rows():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gf = slide.shapes.add_table(4, 2, Emu(500000), Emu(500000), Emu(4000000), Emu(2000000))
    table = gf.table
    # Template ships 4 sample rows.
    for r in range(4):
        for c in range(2):
            table.cell(r, c).text = f"sample{r}{c}"
    # Deck provides only 2 rows of real data.
    _fill_table(gf, [["A0", "A1"], ["B0", "B1"]])
    assert table.cell(0, 0).text == "A0"
    assert table.cell(1, 1).text == "B1"
    # Surplus template rows 2-3 are now blank, not leftover sample data.
    assert all(table.cell(r, c).text == "" for r in (2, 3) for c in range(2))


from pptx_mcp.filler import _grow_table_rows


def test_grow_table_rows_adds_until_needed():
    prs, slide, gf, table = _deck_with_table(2, 3)
    added = _grow_table_rows(table, 5)
    assert added == 3
    assert len(table.rows) == 5
    assert len(table.columns) == 3   # columns unchanged


def test_grow_table_rows_noop_when_enough():
    prs, slide, gf, table = _deck_with_table(4, 2)
    assert _grow_table_rows(table, 4) == 0
    assert _grow_table_rows(table, 2) == 0
    assert len(table.rows) == 4


def test_fill_table_grows_grid_for_extra_rows():
    prs, slide, gf, table = _deck_with_table(2, 2)
    warnings = _fill_table(gf, [["A0", "A1"], ["B0", "B1"],
                                ["C0", "C1"], ["D0", "D1"]])
    assert len(table.rows) == 4
    assert table.cell(2, 0).text == "C0"
    assert table.cell(3, 1).text == "D1"
    assert any(w.code == "table_autogrew" for w in warnings)


def test_fill_table_no_grow_warning_when_it_fits():
    prs, slide, gf, table = _deck_with_table(4, 2)
    warnings = _fill_table(gf, [["A0", "A1"], ["B0", "B1"]])
    assert not any(w.code == "table_autogrew" for w in warnings)


def test_fill_shape_text_truncates_and_warns(sample_template_dir):
    from pptx import Presentation
    from pptx_mcp.filler import fill_shape
    from pptx_mcp.models import Constraints
    prs = Presentation(str(sample_template_dir / "base.pptx"))
    slide = prs.slides[0]
    # Find the title textbox (first one added, which is a textbox with text frame)
    shape = next(s for s in slide.shapes if s.has_text_frame)
    long_text = "This is a long text " * 20  # Repeating words to ensure word boundaries
    warns = fill_shape(slide, shape, "text", long_text,
                       Constraints(max_chars=40), slot_id="title")
    assert any(w.code == "text_truncated" for w in warns)
    assert shape.text_frame.text  # Text should not be empty
    assert len(shape.text_frame.text) <= 40  # Should be constrained by max_chars


def test_fill_shape_table_fills(sample_template_dir):
    from pptx import Presentation
    from pptx_mcp.filler import fill_shape
    from pptx_mcp.models import Constraints
    prs = Presentation(str(sample_template_dir / "base.pptx"))
    slide = prs.slides[2]  # table slide
    table_shape = next(s for s in slide.shapes if s.has_table)
    fill_shape(slide, table_shape, "table", [["A", "B"], ["C", "D"]], Constraints())
    assert table_shape.table.cell(0, 0).text == "A"


def test_fill_shape_unknown_kind_is_noop(sample_template_dir):
    from pptx import Presentation
    from pptx_mcp.filler import fill_shape
    from pptx_mcp.models import Constraints
    prs = Presentation(str(sample_template_dir / "base.pptx"))
    slide = prs.slides[0]
    assert fill_shape(slide, slide.shapes[0], "other", "x", Constraints()) == []
