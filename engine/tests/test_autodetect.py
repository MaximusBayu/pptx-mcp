from pptx import Presentation
from pptx_mcp.autodetect import classify_shape, estimate_max_chars, derive_ids, ShapeAssessment, autodetect


def _assess_file_shapes(path):
    prs = Presentation(path)
    out = {}
    sw, sh = prs.slide_width, prs.slide_height
    for slide in prs.slides:
        for shp in slide.shapes:
            out[shp.shape_id] = classify_shape(shp, sw, sh)
    return out


def test_classifier_separates_slots_from_decoration(labeled_deck):
    path, labels = labeled_deck
    assessed = _assess_file_shapes(path)
    for sid, is_slot in labels.items():
        a = assessed[sid]
        assert a.is_candidate == is_slot, f"shape {sid}: conf={a.confidence}"


def test_confidence_in_range(labeled_deck):
    path, _ = labeled_deck
    for a in _assess_file_shapes(path).values():
        assert 0.0 <= a.confidence <= 1.0


def test_max_chars_scales_with_box_and_font():
    emu = 914400
    big = estimate_max_chars(int(8 * emu), int(1.5 * emu), 40.0)[0]
    small_font = estimate_max_chars(int(8 * emu), int(1.5 * emu), 20.0)[0]
    assert small_font > big
    assert big > 0


def test_max_chars_uses_default_font_when_none():
    emu = 914400
    mc, ml = estimate_max_chars(int(4 * emu), int(2 * emu), None)
    assert mc > 0 and ml >= 1


def _mk(sid, x, y, w, h, conf=0.9, typ="text"):
    return ShapeAssessment(sid, f"TextBox {sid}", typ,
                           {"x": x, "y": y, "w": w, "h": h}, conf, True, 18.0)


def test_derive_ids_hybrid_semantic_then_indexed():
    title = _mk(1, 10, 5, 70, 15)
    subtitle = _mk(2, 10, 22, 50, 6)
    body = _mk(3, 10, 35, 70, 45)
    other = _mk(4, 10, 82, 20, 5)
    ids = derive_ids([title, subtitle, body, other])
    assert ids[1] == "title"
    assert ids[2] == "subtitle"
    assert ids[3] == "body"
    assert ids[4].startswith("text_")


def test_autodetect_shapes_have_suggestions(labeled_deck):
    path, labels = labeled_deck
    data = open(path, "rb").read()
    out = autodetect(data)
    shapes = {s["shape_id"]: s for s in out["slides"][0]["shapes"]}
    for sid, is_slot in labels.items():
        assert shapes[sid]["is_candidate"] == is_slot
    for s in out["slides"][0]["shapes"]:
        if s["is_candidate"]:
            assert s["suggested_id"]
            assert s["suggested_max_chars"] > 0


def test_table_candidate_gets_max_rows_cols(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(8), Inches(3))
    p = tmp_path / "tbl.pptx"
    prs.save(str(p))
    det = autodetect(p.read_bytes())
    tables = [s for s in det["slides"][0]["shapes"] if s["type"] == "table"]
    assert tables, "table shape not detected"
    t = tables[0]
    assert t["is_candidate"] is True
    assert t["suggested_max_rows"] == 3
    assert t["suggested_max_cols"] == 2
    # text-only fields stay zero for a table
    assert t["suggested_max_chars"] == 0


def test_shapes_have_text_and_example(labeled_deck):
    path, _ = labeled_deck
    out = autodetect(open(path, "rb").read())
    shapes = out["slides"][0]["shapes"]
    titles = [s for s in shapes if s["suggested_id"] == "title"]
    assert titles, "no title candidate detected"
    t = titles[0]
    # the title box text in the fixture is "Quarterly Business Review"
    assert "Quarterly Business Review" in t["text"]
    # example is the box's own text (the biggest agent win), non-empty for a tagged text slot
    assert t["suggested_example"]
    assert t["suggested_example"] in t["text"] or t["text"].startswith(t["suggested_example"].rstrip("…"))
    assert t["suggested_description"] == "Slide title"


def test_slot_description_labels():
    from pptx_mcp.autodetect import slot_description
    assert slot_description("title") == "Slide title"
    assert slot_description("subtitle") == "Subtitle"
    assert slot_description("body") == "Body text"
    assert slot_description("table_1") == "Table data"
    assert slot_description("image_2") == "Image"
    assert slot_description("text_3") == "Text"
    assert slot_description("") == "Text"


def test_slide_kind_rules():
    from pptx_mcp.autodetect import slide_kind
    assert slide_kind("", False, 0, 2, True) == "cover"
    assert slide_kind("Agenda for today", False, 1, 4, False) == "agenda"
    assert slide_kind("Executive Summary", False, 2, 3, False) == "summary"
    assert slide_kind("Severity: CRITICAL CWE-89", False, 3, 3, False) == "finding"
    assert slide_kind("just numbers", True, 4, 1, False) == "data"
    assert slide_kind("Thank you!", False, 5, 1, False) == "closing"
    assert slide_kind("Section One", False, 6, 1, False) == "section"
    assert slide_kind("a paragraph of content here", False, 7, 4, False) == "content"


def test_slide_description_lists_slots_and_repeat_hint():
    from pptx_mcp.autodetect import slide_description
    d = slide_description("finding", ["title", "severity", "description"])
    assert "title" in d and "severity" in d and "description" in d
    assert "Repeat per item" in d
    assert "no slots" in slide_description("content", [])


def test_autodetect_slide_has_kind(labeled_deck):
    path, _ = labeled_deck
    out = autodetect(open(path, "rb").read())
    slide = out["slides"][0]
    assert slide["kind"] in {
        "cover", "agenda", "summary", "finding", "data", "closing", "section", "content"
    }
    assert slide["suggested_name"] == slide["kind"]
    assert slide["suggested_description"]


def _finding_like(slide, idx):
    from pptx.util import Inches, Pt
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = f"Finding F{idx}: SQL Injection"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    body.text_frame.text = "Severity: HIGH. CWE-89. Remediation steps here."
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)


def test_repeatable_marks_structural_twins(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank = prs.slide_layouts[6]
    # two structurally identical finding slides + one unique cover-ish slide
    _finding_like(prs.slides.add_slide(blank), 1)
    _finding_like(prs.slides.add_slide(blank), 2)
    cover = prs.slides.add_slide(blank)
    c = cover.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    c.text_frame.text = "RISEStore VAPT Report"
    c.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    p = tmp_path / "rep.pptx"
    prs.save(str(p))

    out = autodetect(p.read_bytes())
    slides = out["slides"]
    assert slides[0]["repeatable"] is True
    assert slides[1]["repeatable"] is True
    assert slides[2]["repeatable"] is False


# Tests for table overlap demoting
from pptx_mcp.autodetect import (
    _rect_overlap_frac, _demote_text_in_tables, ShapeAssessment, TABLE_OVERLAP_TAU,
)


def _assess(shape_id, type_, bbox, is_candidate=True, confidence=0.9):
    return ShapeAssessment(
        shape_id=shape_id, name=f"s{shape_id}", type=type_, bbox_pct=bbox,
        confidence=confidence, is_candidate=is_candidate, font_pt=18.0,
    )


def test_rect_overlap_fully_contained():
    a = {"x": 10, "y": 10, "w": 10, "h": 10}   # inside b
    b = {"x": 0, "y": 0, "w": 100, "h": 100}
    assert _rect_overlap_frac(a, b) == 1.0


def test_rect_overlap_disjoint():
    a = {"x": 0, "y": 0, "w": 10, "h": 10}
    b = {"x": 50, "y": 50, "w": 10, "h": 10}
    assert _rect_overlap_frac(a, b) == 0.0


def test_rect_overlap_half_in():
    a = {"x": 0, "y": 0, "w": 10, "h": 10}      # area 100
    b = {"x": 5, "y": 0, "w": 100, "h": 10}     # covers right half of a
    assert abs(_rect_overlap_frac(a, b) - 0.5) < 1e-9


def test_rect_overlap_zero_area_text():
    a = {"x": 0, "y": 0, "w": 0, "h": 10}
    b = {"x": 0, "y": 0, "w": 100, "h": 100}
    assert _rect_overlap_frac(a, b) == 0.0


def test_demote_text_inside_table():
    table = _assess(1, "table", {"x": 0, "y": 0, "w": 80, "h": 80})
    inside = _assess(2, "text", {"x": 10, "y": 10, "w": 20, "h": 20})
    assessments = [table, inside]
    _demote_text_in_tables(assessments)
    assert inside.is_candidate is False
    assert inside.confidence < TABLE_OVERLAP_TAU  # also below TAU=0.5
    assert table.is_candidate is True  # table never demoted


def test_text_beside_table_stays_candidate():
    table = _assess(1, "table", {"x": 0, "y": 40, "w": 80, "h": 50})
    title = _assess(2, "text", {"x": 0, "y": 0, "w": 80, "h": 20})  # above, no overlap
    _demote_text_in_tables([table, title])
    assert title.is_candidate is True


def test_no_tables_no_change():
    a = _assess(1, "text", {"x": 0, "y": 0, "w": 50, "h": 50})
    _demote_text_in_tables([a])
    assert a.is_candidate is True


def _plain_slide(slide, texts):
    from pptx.util import Inches, Pt
    for i, t in enumerate(texts):
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5 + i * 1.3), Inches(8), Inches(1))
        tb.text_frame.text = t
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(18)


def test_single_content_slide_repeatable_by_kind(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    _plain_slide(prs.slides.add_slide(blank), ["Cover Title"])          # idx 0 -> cover
    _plain_slide(prs.slides.add_slide(blank),                            # idx 1 -> content
                 ["Background", "Goals of the project", "Scope notes", "Other detail"])
    p = tmp_path / "content.pptx"
    prs.save(str(p))
    slides = autodetect(p.read_bytes())["slides"]
    assert slides[1]["kind"] == "content"
    assert slides[1]["repeatable"] is True      # single instance, flagged by kind
    assert slides[0]["repeatable"] is False     # cover is not repeatable


def test_single_finding_slide_repeatable_by_kind(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    _plain_slide(prs.slides.add_slide(blank), ["Cover Title"])          # idx 0 -> cover
    _plain_slide(prs.slides.add_slide(blank),                            # idx 1 -> finding
                 ["Finding F1", "Severity: HIGH", "CWE-89 details", "Remediation steps"])
    p = tmp_path / "finding.pptx"
    prs.save(str(p))
    slides = autodetect(p.read_bytes())["slides"]
    assert slides[1]["kind"] == "finding"
    assert slides[1]["repeatable"] is True


def test_single_summary_slide_not_repeatable(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    _plain_slide(prs.slides.add_slide(blank), ["Cover Title"])          # idx 0 -> cover
    _plain_slide(prs.slides.add_slide(blank),                            # idx 1 -> summary
                 ["Executive Summary", "Point one here", "Point two here", "Point three"])
    p = tmp_path / "summary.pptx"
    prs.save(str(p))
    slides = autodetect(p.read_bytes())["slides"]
    assert slides[1]["kind"] == "summary"
    assert slides[1]["repeatable"] is False     # summary kind is not repeatable


def test_slide_description_repeat_hint_for_content():
    from pptx_mcp.autodetect import slide_description
    assert "Repeat per item" in slide_description("content", ["body"])
    assert "Repeat per item" in slide_description("finding", ["title"])
    assert "Repeat per item" not in slide_description("cover", ["title"])
